# -*- coding: utf-8 -*-
"""
全彩版PPT制作工作流 V1.1.0
========================
每张页面都是AI生成的唯美背景图
文字内容直接排版，支持可视化Markdown格式
带图形图示的精美PPT

V1.1.0 新增：
- ⭐ 重要提示框
- 📊 数据标签
- ✓ 待办/完成标记
- 💡 技巧提示框
- ⚠️ 注意警告框
- 🔢 数字编号
- 📈 进度指示器
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# 配色方案
COLORS = {
    'title_bar': RGBColor(255, 255, 255),      # 白色标题条
    'h2_title': RGBColor(100, 60, 120),        # 二级标题紫色
    'h3_title': RGBColor(120, 80, 150),        # 三级标题
    'h2_decorator': RGBColor(200, 100, 150),   # 装饰线粉色
    'bullet': RGBColor(60, 60, 90),            # 列表文字
    'bold': RGBColor(100, 50, 120),            # 粗体
    'divider': RGBColor(200, 180, 220),        # 分隔线
    
    # V1.1.0 新增配色
    'tip_bg': RGBColor(255, 248, 220),         # 提示框背景-米色
    'tip_border': RGBColor(255, 200, 100),     # 提示框边框-金色
    'warning_bg': RGBColor(255, 240, 240),     # 警告框背景-淡红
    'warning_border': RGBColor(255, 100, 100),  # 警告框边框-红色
    'important_bg': RGBColor(255, 245, 250),    # 重要提示背景-淡粉
    'important_border': RGBColor(200, 80, 120),  # 重要提示边框-深粉
    'success_bg': RGBColor(240, 255, 240),      # 成功框背景-淡绿
    'success_border': RGBColor(100, 180, 100),  # 成功框边框-绿色
    'tag_bg': RGBColor(230, 220, 250),          # 标签背景-淡紫
    'number_bg': RGBColor(150, 100, 180),       # 数字编号背景-紫色
    'progress_bg': RGBColor(220, 220, 240),     # 进度条背景
    'progress_fill': RGBColor(180, 120, 200),    # 进度条填充
    'quote_bg': RGBColor(245, 240, 255),        # 引用框背景
}


def add_background_image(slide, image_path):
    """添加全屏背景图"""
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path,
            Inches(0), Inches(0),
            width=Inches(13.333),
            height=Inches(7.5)
        )


def add_title_bar(slide, text, left=0.3, top=0.2, width=12.7, height=0.7):
    """添加白色标题条"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['title_bar']
    bar.line.fill.background()
    bar.adjustments[0] = 0.05
    
    tb = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.12),
        Inches(width - 0.4), Inches(0.5)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 60, 130)


def add_h2_title(slide, text, left=0.4, top=1.1, width=12.5):
    """二级标题 - 带粉色装饰线"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(0.6), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['h2_decorator']
    line.line.fill.background()
    
    tb = slide.shapes.add_textbox(
        Inches(left + 0.7), Inches(top - 0.05),
        Inches(width - 0.7), Inches(0.5)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['h2_title']


def add_h3_title(slide, text, left=0.5, top=1.6, width=12):
    """三级标题 - 带◈标记"""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(0.4)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "◈ " + text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['h3_title']


def add_bullet(slide, text, left=0.6, top=2.1, width=12, font_size=14):
    """列表项"""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(0.35)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    
    if '**' in text:
        parts = text.split('**')
        if len(parts) >= 3:
            run1 = p
            run1.text = "    ●  " + parts[0]
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = COLORS['bullet']
            
            run2 = tf.add_paragraph()
            run2.text = parts[1]
            run2.font.size = Pt(font_size + 1)
            run2.font.bold = True
            run2.font.color.rgb = COLORS['bold']
            
            if len(parts) > 2 and parts[2]:
                run3 = tf.add_paragraph()
                run3.text = parts[2]
                run3.font.size = Pt(font_size)
                run3.font.color.rgb = COLORS['bullet']
            return
    
    p.text = "    ●  " + text
    p.font.size = Pt(font_size)
    p.font.color.rgb = COLORS['bullet']


def add_divider(slide, left=0.5, top=6.8, width=12.3):
    """分隔线"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['divider']
    line.line.fill.background()


# ==================== V1.1.0 新增图形图示 ====================

def add_tip_box(slide, text, left=0.4, top=2.0, width=12.5, icon="💡"):
    """💡 技巧提示框"""
    # 背景框
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.7)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['tip_bg']
    box.line.color.rgb = COLORS['tip_border']
    box.line.width = Pt(1.5)
    box.adjustments[0] = 0.1
    
    # 图标+文字
    tb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.12),
        Inches(width - 0.3), Inches(0.5)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = icon + " "
    run1.font.size = Pt(16)
    run1.font.color.rgb = RGBColor(200, 150, 50)
    
    run2 = p.add_run()
    run2.text = text
    run2.font.size = Pt(13)
    run2.font.color.rgb = RGBColor(120, 90, 60)
    
    return box


def add_warning_box(slide, text, left=0.4, top=2.0, width=12.5):
    """⚠️ 注意警告框"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.7)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['warning_bg']
    box.line.color.rgb = COLORS['warning_border']
    box.line.width = Pt(1.5)
    box.adjustments[0] = 0.1
    
    tb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.12),
        Inches(width - 0.3), Inches(0.5)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "⚠️ "
    run1.font.size = Pt(16)
    run1.font.color.rgb = RGBColor(200, 50, 50)
    
    run2 = p.add_run()
    run2.text = text
    run2.font.size = Pt(13)
    run2.font.color.rgb = RGBColor(150, 50, 50)
    
    return box


def add_important_box(slide, text, left=0.4, top=2.0, width=12.5):
    """⭐ 重要提示框"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.7)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['important_bg']
    box.line.color.rgb = COLORS['important_border']
    box.line.width = Pt(2)
    box.adjustments[0] = 0.1
    
    tb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.12),
        Inches(width - 0.3), Inches(0.5)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "⭐ "
    run1.font.size = Pt(16)
    run1.font.color.rgb = RGBColor(200, 100, 50)
    
    run2 = p.add_run()
    run2.text = text
    run2.font.size = Pt(13)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(150, 50, 80)
    
    return box


def add_success_box(slide, text, left=0.4, top=2.0, width=12.5):
    """✓ 成功/完成框"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.65)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['success_bg']
    box.line.color.rgb = COLORS['success_border']
    box.line.width = Pt(1.5)
    box.adjustments[0] = 0.1
    
    tb = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.1),
        Inches(width - 0.3), Inches(0.45)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "✓ "
    run1.font.size = Pt(15)
    run1.font.color.rgb = RGBColor(50, 150, 50)
    run1.font.bold = True
    
    run2 = p.add_run()
    run2.text = text
    run2.font.size = Pt(12)
    run2.font.color.rgb = RGBColor(50, 120, 50)
    
    return box


def add_numbered_item(slide, number, text, left=0.5, top=2.0, width=12.3):
    """🔢 数字编号项"""
    # 数字圆圈
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top + 0.02), Inches(0.4), Inches(0.4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['number_bg']
    circle.line.fill.background()
    
    # 数字
    tb_num = slide.shapes.add_textbox(
        Inches(left + 0.05), Inches(top + 0.05),
        Inches(0.3), Inches(0.35)
    )
    tf_num = tb_num.text_frame
    p_num = tf_num.paragraphs[0]
    p_num.text = str(number)
    p_num.font.size = Pt(14)
    p_num.font.bold = True
    p_num.font.color.rgb = RGBColor(255, 255, 255)
    p_num.alignment = PP_ALIGN.CENTER
    
    # 内容
    tb = slide.shapes.add_textbox(
        Inches(left + 0.5), Inches(top + 0.05),
        Inches(width - 0.5), Inches(0.35)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = COLORS['bullet']
    
    return circle


def add_tag(slide, text, left=0.5, top=2.0):
    """🏷️ 标签"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(len(text) * 0.15 + 0.4), Inches(0.35)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['tag_bg']
    box.line.fill.background()
    box.adjustments[0] = 0.5  # 高度圆角
    
    tb = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.05),
        Inches(len(text) * 0.15 + 0.2), Inches(0.3)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(100, 60, 120)
    
    return box


def add_progress_bar(slide, percent, left=0.5, top=2.0, width=3.0, label=""):
    """📈 进度指示器"""
    # 标签
    if label:
        tb_label = slide.shapes.add_textbox(
            Inches(left), Inches(top - 0.25), Inches(1), Inches(0.25)
        )
        tf_label = tb_label.text_frame
        p_label = tf_label.paragraphs[0]
        p_label.text = label
        p_label.font.size = Pt(10)
        p_label.font.color.rgb = RGBColor(100, 100, 130)
    
    # 背景条
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.2)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['progress_bg']
    bg.line.fill.background()
    bg.adjustments[0] = 0.5
    
    # 填充条
    fill_width = width * (percent / 100)
    if fill_width > 0.1:
        fill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(fill_width), Inches(0.2)
        )
        fill.fill.solid()
        fill.fill.fore_color.rgb = COLORS['progress_fill']
        fill.line.fill.background()
        fill.adjustments[0] = 0.5
    
    # 百分比
    tb_pct = slide.shapes.add_textbox(
        Inches(left + width + 0.1), Inches(top - 0.02),
        Inches(0.6), Inches(0.25)
    )
    tf_pct = tb_pct.text_frame
    p_pct = tf_pct.paragraphs[0]
    p_pct.text = f"{percent}%"
    p_pct.font.size = Pt(10)
    p_pct.font.bold = True
    p_pct.font.color.rgb = RGBColor(150, 100, 160)
    
    return bg


def add_quote_box(slide, text, left=0.5, top=2.0, width=12.3):
    """"引用框"""
    # 左边竖线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(0.08), Inches(0.8)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['number_bg']
    line.line.fill.background()
    
    # 背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + 0.08), Inches(top), Inches(width - 0.08), Inches(0.8)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['quote_bg']
    bg.line.fill.background()
    
    # 文字
    tb = slide.shapes.add_textbox(
        Inches(left + 0.25), Inches(top + 0.15),
        Inches(width - 0.4), Inches(0.6)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"' + text + '"'
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 80, 130)
    
    return line


def add_emoji_icon(slide, emoji, left=0.4, top=2.0):
    """添加表情图标"""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(0.5), Inches(0.5)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(22)
    return tb


# ==================== 页面类型 ====================

def add_cover_page(slide, title, subtitle=None):
    """封面页"""
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(2.5), Inches(9.333), Inches(2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS['title_bar']
    title_bar.line.fill.background()
    title_bar.adjustments[0] = 0.1
    
    tb = slide.shapes.add_textbox(
        Inches(2.5), Inches(2.8), Inches(8.333), Inches(1)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 60, 130)
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle and subtitle != title:
        tb2 = slide.shapes.add_textbox(
            Inches(3), Inches(4.0), Inches(7.333), Inches(0.5)
        )
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(120, 120, 160)
        p2.alignment = PP_ALIGN.CENTER


def add_section_page(slide, title):
    """章节分隔页"""
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(2.8), Inches(10.333), Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLORS['title_bar']
    title_bar.line.fill.background()
    title_bar.adjustments[0] = 0.08
    
    tb = slide.shapes.add_textbox(
        Inches(2), Inches(3.0), Inches(9.333), Inches(0.8)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 60, 130)
    p.alignment = PP_ALIGN.CENTER


# ==================== 主函数 ====================

def parse_special_blocks(text):
    """解析特殊格式块"""
    lines = text.split('\n')
    blocks = []
    current_block = None
    current_content = []
    
    for line in lines:
        line = line.strip()
        
        # 检测特殊块开始
        if line.startswith(':::tip ') or line.startswith(':::tip'):
            if current_block:
                blocks.append(current_block)
            current_block = {'type': 'tip', 'content': line.replace(':::tip', '').strip()}
        elif line.startswith(':::warning ') or line.startswith(':::warning'):
            if current_block:
                blocks.append(current_block)
            current_block = {'type': 'warning', 'content': line.replace(':::warning', '').strip()}
        elif line.startswith(':::important ') or line.startswith(':::important'):
            if current_block:
                blocks.append(current_block)
            current_block = {'type': 'important', 'content': line.replace(':::important', '').strip()}
        elif line.startswith(':::success ') or line.startswith(':::success'):
            if current_block:
                blocks.append(current_block)
            current_block = {'type': 'success', 'content': line.replace(':::success', '').strip()}
        elif line == ':::':
            if current_block:
                blocks.append(current_block)
                current_block = None
        elif current_block:
            current_block['content'] += '\n' + line
        else:
            blocks.append({'type': 'normal', 'content': line})
    
    if current_block:
        blocks.append(current_block)
    
    return blocks


def create_fullcolor_ppt(slides, output_path, background_images=None):
    """
    创建全彩版PPT V1.1.0
    
    新增特殊块支持：
    - :::tip 技巧提示
    - :::warning 注意警告
    - :::important 重要提示
    - :::success 成功完成
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    
    bg_idx = 0
    bg_count = len(background_images) if background_images else 0
    
    for idx, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        
        if bg_count > 0:
            add_background_image(slide, background_images[bg_idx % bg_count])
            bg_idx += 1
        
        stype = slide_data.get('type', '')
        title = slide_data.get('title', '') or slide_data.get('main', '')
        items = slide_data.get('items', [])
        
        if stype == 'title':
            add_cover_page(slide, title, slide_data.get('sub'))
        
        elif stype == 'section':
            add_section_page(slide, title)
        
        elif stype in ['content_list', 'content_cards']:
            add_title_bar(slide, title)
            
            y_pos = 1.1
            max_y = 7.0
            
            for item in items:
                if y_pos > max_y:
                    break
                
                # 特殊块处理
                if isinstance(item, str) and item.startswith(':::tip'):
                    content = item.replace(':::tip', '').strip()
                    if content:
                        add_tip_box(slide, content, top=y_pos)
                        y_pos += 0.85
                    continue
                elif isinstance(item, str) and item.startswith(':::warning'):
                    content = item.replace(':::warning', '').strip()
                    if content:
                        add_warning_box(slide, content, top=y_pos)
                        y_pos += 0.85
                    continue
                elif isinstance(item, str) and item.startswith(':::important'):
                    content = item.replace(':::important', '').strip()
                    if content:
                        add_important_box(slide, content, top=y_pos)
                        y_pos += 0.85
                    continue
                elif isinstance(item, str) and item.startswith(':::success'):
                    content = item.replace(':::success', '').strip()
                    if content:
                        add_success_box(slide, content, top=y_pos)
                        y_pos += 0.8
                    continue
                
                if isinstance(item, dict):
                    t = item.get('title', '')
                    v = item.get('value', '')
                    is_group = item.get('is_group_header', False)
                    
                    if is_group and t:
                        add_h2_title(slide, t, top=y_pos)
                        y_pos += 0.55
                    elif t:
                        add_h3_title(slide, t, top=y_pos)
                        y_pos += 0.4
                        if v and y_pos < max_y:
                            add_bullet(slide, v, top=y_pos)
                            y_pos += 0.4
                        y_pos += 0.1
                    elif v:
                        add_bullet(slide, v, top=y_pos)
                        y_pos += 0.4
                        
                elif isinstance(item, str):
                    item = item.strip()
                    if not item:
                        y_pos += 0.15
                        continue
                    if item.startswith('## '):
                        add_h2_title(slide, item[3:], top=y_pos)
                        y_pos += 0.55
                    elif item.startswith('### '):
                        add_h3_title(slide, item[4:], top=y_pos)
                        y_pos += 0.45
                    else:
                        add_bullet(slide, item, top=y_pos)
                        y_pos += 0.4
            
            add_divider(slide, top=y_pos + 0.2)
        
        else:
            if title:
                add_title_bar(slide, title)
    
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path


if __name__ == '__main__':
    print("全彩版PPT制作工作流 V1.1.0")
    print("新增：💡提示框 ⚠️警告框 ⭐重要提示 ✓完成框 🔢编号 📈进度条")
