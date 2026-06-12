#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT生成器 - V3.0（美学升级版）
核心原则：层次感 / 留白 / 图形化 / 节奏 / 网格 / 克制
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
# ══════════════════════════════════════════════════════════════
# 排版常量 - 网格系统
# ══════════════════════════════════════════════════════════════
class Layout:
    """排版网格系统"""
    # 页面尺寸 16:9
    WIDTH = Inches(13.333)
    HEIGHT = Inches(7.5)
    
    # 边距系统
    MARGIN_LEFT = Inches(0.8)
    MARGIN_RIGHT = Inches(0.8)
    MARGIN_TOP = Inches(0.6)
    MARGIN_BOTTOM = Inches(0.6)
    
    # 内容区
    CONTENT_WIDTH = WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    CONTENT_HEIGHT = HEIGHT - MARGIN_TOP - MARGIN_BOTTOM
# ══════════════════════════════════════════════════════════════
# 主题配色
# ══════════════════════════════════════════════════════════════
THEMES = {
    "executive": {
        "bg": RGBColor(13, 27, 42),        # 深蓝背景
        "card_bg": RGBColor(27, 40, 56),    # 卡片背景
        "text": RGBColor(255, 255, 255),    # 白色文字
        "muted": RGBColor(136, 153, 170),   # 灰色副文字
        "accent": RGBColor(255, 183, 0),    # 金色强调
        "light_accent": RGBColor(255, 220, 100),  # 浅金色
        "success": RGBColor(0, 230, 118),   # 绿色（成功）
        "warning": RGBColor(255, 152, 0),   # 橙色（警告）
        "danger": RGBColor(244, 67, 54),    # 红色（危险）
        "border": RGBColor(45, 63, 82),     # 边框色
    },
    "clean": {  # 白色主题 - 用于AI配图PPT
        "bg": RGBColor(255, 255, 255),      # 纯白背景
        "card_bg": RGBColor(248, 248, 248), # 浅灰卡片
        "text": RGBColor(33, 33, 33),      # 深灰文字
        "muted": RGBColor(120, 120, 120),  # 灰色副文字
        "accent": RGBColor(59, 130, 246),   # 蓝色强调
        "light_accent": RGBColor(96, 165, 250),  # 浅蓝色
        "success": RGBColor(34, 197, 94),   # 绿色
        "warning": RGBColor(245, 158, 11),  # 橙色
        "danger": RGBColor(239, 68, 68),    # 红色
        "border": RGBColor(229, 229, 229),  # 边框色
    },
    "spark": {
        "bg": RGBColor(15, 12, 41),
        "card_bg": RGBColor(30, 25, 60),
        "text": RGBColor(240, 240, 240),
        "muted": RGBColor(160, 176, 192),
        "accent": RGBColor(36, 198, 220),
        "light_accent": RGBColor(100, 220, 240),
        "border": RGBColor(60, 50, 100),
    },
    "dark": {
        "bg": RGBColor(26, 26, 26),
        "card_bg": RGBColor(40, 40, 40),
        "text": RGBColor(255, 255, 255),
        "muted": RGBColor(150, 150, 150),
        "accent": RGBColor(0, 200, 100),
        "light_accent": RGBColor(100, 255, 150),
        "border": RGBColor(60, 60, 60),
    },
}
DEFAULT_THEME = "executive"
class PPTGeneratorV3:
    """V3.0 PPT生成器 - 美学升级版"""
    
    def __init__(self, theme=DEFAULT_THEME):
        self.prs = Presentation()
        self.prs.slide_width = Layout.WIDTH
        self.prs.slide_height = Layout.HEIGHT
        self.t = THEMES.get(theme, THEMES[DEFAULT_THEME])
    
    # ══════════════════════════════════════════════════════════
    # 核心方法：建立视觉层次
    # ══════════════════════════════════════════════════════════
    
    def _add_bg(self, slide):
        """添加背景"""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Layout.WIDTH, Layout.HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.t['bg']
        bg.line.fill.background()
        # 移到最底层
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
    
    def _add_decoration_image(self, slide_idx, image_path):
        """添加装饰图片到页面"""
        # 找到当前slide
        slide = self.prs.slides[slide_idx]
        # 默认右侧装饰
        pic = slide.shapes.add_picture(
            image_path,
            Inches(10.5),  # 右侧
            Inches(0.3),
            width=Inches(2.5)
        )
        # 移到背景层上方
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(3, pic._element)
    
    def _add_image_ex(self, slide, image_path, position="right", size="large", opacity=1.0, add_overlay=False):
        """
        精确控制配图的位置和大小
        
        Args:
            slide: 幻灯片对象
            image_path: 图片路径
            position: 位置 (background/right/right_large/left/etc)
            size: 尺寸 (full/large/medium/small/tiny)
            opacity: 透明度 (0-1)
            add_overlay: 是否添加遮罩
        """
        if not image_path or not os.path.exists(image_path):
            return None
        
        W = 13.333  # 页面宽度 inch
        H = 7.5     # 页面高度 inch
        
        # 尺寸映射（相对于页面宽度）
        size_map = {
            "full": 1.0,    # 100% 页面宽度
            "large": 0.6,   # 60%
            "medium": 0.5,  # 50%
            "small": 0.4,   # 40%
            "tiny": 0.3    # 30%
        }
        width_ratio = size_map.get(size, 0.5)
        
        # 位置映射 - 右侧装饰：图片占右侧区域
        # 右侧60%宽度：x=40%，width=60%
        # 右侧50%宽度：x=50%，width=50%
        W = 13.333  # 页面宽度 inch
        H = 7.5     # 页面高度 inch
        
        # 左右分栏布局（内容55% + 配图45%，中间留间隙）
        # 页面宽度13.33in，边距0.8in，内容区11.73in
        # 内容区：左0.8in-7.25in(6.45in)
        # 间隙：0.1in
        # 配图区：7.35in-12.63in(5.28in)
        position_map = {
            "background": (0, 0, W * 1.0, H * 1.0, True),    # 全屏背景
            "left": (0.8, H * 0.15, 6.45, H * 0.7, True),  # 左侧内容区(55%)
            "right_image": (7.35, H * 0.15, 5.28, H * 0.7, True),  # 右侧配图(45%)
            "right": (W * 0.4, 0, W * 0.6, H * 1.0, True),         # 右侧60%全高
            "right_large": (W * 0.35, 0, W * 0.65, H * 1.0, True), # 右侧65%全高
            "right_medium": (W * 0.5, H * 0.15, W * 0.45, H * 0.7, True),  # 右侧45%，不挡标题
            "decorative": (W * 0.55, H * 0.15, W * 0.4, H * 0.7, True),  # 右侧装饰
            "decorative_right": (W * 0.6, H * 0.2, W * 0.35, H * 0.6, True),  # 右侧小装饰
            "right_small": (W * 0.65, H * 0.2, W * 0.30, H * 0.6, True),  # 右侧小图(30%)
            "right_tiny": (W * 0.7, H * 0.25, W * 0.25, H * 0.5, True),  # 右侧微型(25%)
            "top_right": (W * 0.6, 0, W * 0.4, H * 1.0, True),     # 右上40%
            "bottom_right": (W * 0.6, H * 0.5, W * 0.4, H * 0.5, True),  # 右下40%
            "header": (0, 0, W * 1.0, H * 0.15, True),
            "wide_banner": (0, 0, W * 1.0, H * 0.35, True),
        }
        
        x, y, width, height, keep_aspect = position_map.get(position, (W * 0.4, 0, W * 0.6, H, True))
        
        # 添加图片 - 保持原始宽高比，不拉伸
        if keep_aspect:
            # 只设置宽度，高度自动按比例计算
            pic = slide.shapes.add_picture(
                image_path,
                Inches(x), Inches(y),
                width=Inches(width)
            )
        else:
            # 同时设置宽高
            pic = slide.shapes.add_picture(
                image_path,
                Inches(x), Inches(y),
                width=Inches(width),
                height=Inches(height)
            )
        
        # 移到最上层
        spTree = slide.shapes._spTree
        sp = pic._element
        spTree.remove(sp)
        spTree.append(sp)  # 添加到最末尾（最上层）
        
        # 如果需要遮罩（在图片下方）
        if add_overlay:
            overlay = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = self.t['bg']
            overlay.line.fill.background()
            spTree = slide.shapes._spTree
            sp2 = overlay._element
            spTree.remove(sp2)
            spTree.insert(len(spTree) - 1, sp2)  # 在图片下方
        
        return pic
    
    def _add_background_image(self, slide_idx, image_path):
        """添加背景图片"""
        slide = self.prs.slides[slide_idx]
        # 全屏背景图
        pic = slide.shapes.add_picture(
            image_path,
            0, 0,
            width=Layout.WIDTH
        )
        # 移到最底层
        spTree = slide.shapes._spTree
        sp = pic._element
        spTree.remove(sp)
        spTree.insert(2, sp)
    
        return slide

    def page_title_with_image(self, main, sub=None, tag=None, image_path=None):
        """带AI配图的封面页 - 无全屏遮罩，让图片完全显示"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        t = self.t
        
        spTree = slide.shapes._spTree
        
        # 1. 背景色
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Layout.WIDTH, Layout.HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = t['bg']
        bg.line.fill.background()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)
        
        # 2. 全屏背景图片（保持原始宽高比，不拉伸，16:9适配）
        # PPT尺寸是13.333x7.5英寸（16:9比例）
        if image_path and os.path.exists(image_path):
            from PIL import Image as PILImage
            try:
                # 读取图片原始尺寸
                with PILImage.open(image_path) as img:
                    orig_width, orig_height = img.size
                    img_aspect = orig_width / orig_height
                    page_aspect = 16 / 9  # 13.333 / 7.5
                    
                    # 计算等比缩放后的尺寸
                    if img_aspect > page_aspect:
                        # 图片更扁：以高度为基准，宽度超出
                        target_height = Layout.HEIGHT
                        target_width = int(target_height * img_aspect)
                    else:
                        # 图片更高：以宽度为基准，高度超出
                        target_width = Layout.WIDTH
                        target_height = int(target_width / img_aspect)
                    
                    # 添加图片（居中放置，超出部分会被裁剪）
                    pic = slide.shapes.add_picture(
                        image_path,
                        int((Layout.WIDTH - target_width) / 2),
                        int((Layout.HEIGHT - target_height) / 2),
                        width=target_width,
                        height=target_height
                    )
            except ImportError:
                # 如果没有PIL，用简单方式：只设置宽度，让高度自动按比例
                pic = slide.shapes.add_picture(
                    image_path,
                    0, 0,
                    width=Layout.WIDTH
                )
                # 垂直居中
                if pic.height < Layout.HEIGHT:
                    pic.top = int((Layout.HEIGHT - pic.height) / 2)
            
            # 移到背景层
            spTree.remove(pic._element)
            spTree.insert(3, pic._element)
        
        # 3. V4.9.3修复：标题加底色块+金色描边，确保任何背景都清晰可读
        # 计算色块位置（居中偏下）
        title_bg_left = Layout.MARGIN_LEFT
        title_bg_top = Inches(2.2)
        title_bg_width = Inches(10.5)
        title_bg_height = Inches(3.5)
        
        # 标题底色块
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            title_bg_left, title_bg_top,
            title_bg_width, title_bg_height
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色底色块
        title_bg.line.fill.background()  # 先无边框
        
        # 移到文字下方、图片上方
        spTree.remove(title_bg._element)
        spTree.insert(4, title_bg._element)
        
        # 金色描边框（用另一个矩形实现）
        border = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            title_bg_left - Inches(0.05), title_bg_top - Inches(0.05),
            title_bg_width + Inches(0.1), title_bg_height + Inches(0.1)
        )
        border.fill.background()  # 无填充
        border.line.color.rgb = t.get('accent', RGBColor(255, 183, 0))  # 金色边框
        border.line.width = Pt(3)
        spTree.remove(border._element)
        spTree.insert(5, border._element)
        
        # 主标题（白色底色块上用黑色文字）
        title_box = slide.shapes.add_textbox(
            title_bg_left + Inches(0.5), title_bg_top + Inches(0.4),
            title_bg_width - Inches(1), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = main
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)  # 白色底色块上用黑色
        
        if sub:
            sub_box = slide.shapes.add_textbox(
                title_bg_left + Inches(0.5), title_bg_top + Inches(2.0),
                title_bg_width - Inches(1), Inches(0.8)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = sub
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(80, 80, 80)  # 深灰色
        
        if tag:
            tag_box = slide.shapes.add_textbox(
                title_bg_left + Inches(0.5), title_bg_top + Inches(2.8),
                title_bg_width - Inches(1), Inches(0.5)
            )
            tf = tag_box.text_frame
            p = tf.paragraphs[0]
            p.text = tag
            p.font.size = Pt(16)
            p.font.color.rgb = t.get('accent', RGBColor(255, 183, 0))
        
        return slide
    
    def page_big_number_with_image(self, number, label, context=None, image_path=None):
        """带AI配图的大数据页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 添加装饰图
        if image_path and os.path.exists(image_path):
            pic = slide.shapes.add_picture(
                image_path,
                Inches(9.5), Inches(1),
                width=Inches(3.5)
            )
            spTree = slide.shapes._spTree
            sp = pic._element
            spTree.remove(sp)
            spTree.insert(3, sp)
        
        # 数字
        num_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Inches(1.5),
            Inches(8), Inches(3.5)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(160)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        
        # 标签
        label_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Inches(5.2),
            Inches(8), Inches(0.8)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(36)
        p.font.color.rgb = t['text']
        
        # 上下文
        if context:
            ctx_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(6.2),
                Inches(8), Inches(0.5)
            )
            tf = ctx_box.text_frame
            p = tf.paragraphs[0]
            p.text = context
            p.font.size = Pt(20)
            p.font.color.rgb = t['muted']
        
        return slide
    
    def page_closing_with_image(self, title, cta=None, contact=None, image_path=None):
        """带AI配图的结尾页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 添加AI配图（16:9适配，保持原始宽高比，不拉伸）
        if image_path and os.path.exists(image_path):
            from PIL import Image as PILImage
            try:
                # 读取图片原始尺寸
                with PILImage.open(image_path) as img:
                    orig_width, orig_height = img.size
                    img_aspect = orig_width / orig_height
                    page_aspect = 16 / 9
                    
                    # 计算等比缩放后的尺寸
                    if img_aspect > page_aspect:
                        target_height = Layout.HEIGHT
                        target_width = int(target_height * img_aspect)
                    else:
                        target_width = Layout.WIDTH
                        target_height = int(target_width / img_aspect)
                    
                    # 添加图片（居中放置，超出部分会被裁剪）
                    pic = slide.shapes.add_picture(
                        image_path,
                        int((Layout.WIDTH - target_width) / 2),
                        int((Layout.HEIGHT - target_height) / 2),
                        width=target_width,
                        height=target_height
                    )
            except ImportError:
                pic = slide.shapes.add_picture(
                    image_path,
                    Inches(0), Inches(0),
                    width=Layout.WIDTH
                )
                if pic.height < Layout.HEIGHT:
                    pic.top = int((Layout.HEIGHT - pic.height) / 2)
            
            # 移到背景
            spTree = slide.shapes._spTree
            sp = pic._element
            spTree.remove(sp)
            spTree.insert(2, sp)
            
            # 无全屏遮罩，让图片显示
        
        # 主标题
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Inches(2.5),
            Layout.CONTENT_WIDTH, Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = t['text']
        p.alignment = PP_ALIGN.CENTER
        
        # CTA
        if cta:
            cta_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(4.5),
                Layout.CONTENT_WIDTH, Inches(0.8)
            )
            tf = cta_box.text_frame
            p = tf.paragraphs[0]
            p.text = cta
            p.font.size = Pt(24)
            p.font.color.rgb = t['accent']
            p.alignment = PP_ALIGN.CENTER
        
        # 联系
        if contact:
            contact_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(6),
                Layout.CONTENT_WIDTH, Inches(0.5)
            )
            tf = contact_box.text_frame
            p = tf.paragraphs[0]
            p.text = contact
            p.font.size = Pt(16)
            p.font.color.rgb = t['muted']
            p.alignment = PP_ALIGN.CENTER
        
        return slide
        """超大字标题 - 核心视觉元素"""
        box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, 
            Inches(2.0), 
            Layout.CONTENT_WIDTH, 
            Inches(3)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = color or self.t['text']
        p.alignment = PP_ALIGN.LEFT
        return box
    
    def _subtitle_small(self, slide, text, top=Inches(4.5)):
        """小字副标题"""
        box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT,
            top,
            Layout.CONTENT_WIDTH,
            Inches(1)
        )
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.color.rgb = self.t['muted']
        p.alignment = PP_ALIGN.LEFT
        return box
    
    def _accent_line(self, slide, top=Inches(1.8)):
        """强调线 - 视觉引导"""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Layout.MARGIN_LEFT, top,
            Inches(1.5), Inches(0.08)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.t['accent']
        line.line.fill.background()
        return line
    
    def _card(self, slide, left, top, width, height):
        """卡片容器"""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.t['card_bg']
        card.line.fill.background()
        return card
    
    # ══════════════════════════════════════════════════════════
    # 页面类型
    # ══════════════════════════════════════════════════════════
    
    def page_title(self, main, sub=None, tag=None):
        """🎯 标题页 - 大留白，大标题"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 左侧强调线
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Layout.MARGIN_LEFT, Inches(2.5),
            Inches(0.12), Inches(2.5)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = t['accent']
        accent.line.fill.background()
        
        # 主标题 - 超大！
        title_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(2.3), 
            Layout.CONTENT_WIDTH - Inches(0.4), Inches(2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = main
        p.font.size = Pt(66)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 副标题
        if sub:
            sub_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(4.5),
                Layout.CONTENT_WIDTH, Inches(1)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = sub
            p.font.size = Pt(28)
            p.font.color.rgb = t['muted']
        
        # 标签（右下角）
        if tag:
            tag_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(6.2),
                Layout.CONTENT_WIDTH, Inches(0.5)
            )
            tf = tag_box.text_frame
            p = tf.paragraphs[0]
            p.text = tag
            p.font.size = Pt(16)
            p.font.color.rgb = t['accent']
        
        return slide
    
    def page_big_number(self, number, label, context=None):
        """📊 大数据页 - 数字震撼"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 数字 - 超级大！
        num_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Inches(1.5),
            Layout.CONTENT_WIDTH, Inches(3.5)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(160)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        
        # 标签 - 中等
        label_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Inches(5.2),
            Layout.CONTENT_WIDTH, Inches(0.8)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(36)
        p.font.color.rgb = t['text']
        
        # 上下文 - 小
        if context:
            ctx_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(6.2),
                Layout.CONTENT_WIDTH, Inches(0.5)
            )
            tf = ctx_box.text_frame
            p = tf.paragraphs[0]
            p.text = context
            p.font.size = Pt(20)
            p.font.color.rgb = t['muted']
        
        return slide
    
    def page_two_columns(self, left_title, left_items, right_title, right_items, page_title=None):
        """📐 对比页 - 左右分栏"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 页面标题
        if page_title:
            title_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
                Layout.CONTENT_WIDTH, Inches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = page_title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = t['text']
        
        content_top = Inches(1.8) if page_title else Inches(1.2)
        
        # 左栏
        left_width = (Layout.CONTENT_WIDTH - Inches(0.4)) / 2
        
        # 左栏标题
        lt_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, content_top,
            left_width, Inches(0.6)
        )
        tf = lt_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        
        # 左栏内容 - 大留白
        lc_top = content_top + Inches(0.8)
        lc_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, lc_top,
            left_width, Inches(4.5)
        )
        tf = lc_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = t['text']
            p.space_before = Pt(18)
        
        # 中间分隔线
        mid_x = Layout.MARGIN_LEFT + left_width + Inches(0.2)
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, mid_x, content_top, Inches(0.02), Inches(5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = t['border']
        divider.line.fill.background()
        
        # 右栏
        right_left = mid_x + Inches(0.2)
        
        rt_box = slide.shapes.add_textbox(
            right_left, content_top,
            left_width, Inches(0.6)
        )
        tf = rt_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        
        rc_box = slide.shapes.add_textbox(
            right_left, lc_top,
            left_width, Inches(4.5)
        )
        tf = rc_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = t['text']
            p.space_before = Pt(18)
        
        return slide
    
    def page_chart_visual(self, data, chart_type='radar', has_side_image=False, chart_theme='clean'):
        """📈 图表页 - 图表+文字配合，左右分栏
        
        Args:
            has_side_image: 是否有右侧配图，如有则图表区域缩小
            chart_theme: 图表配色主题，默认白色主题
        """
        from chart_generator import (
            create_radar_chart, create_funnel_chart,
            create_matrix_chart, create_pyramid_chart,
            create_gantt_chart, create_heatmap_chart,
            create_gauge_chart, create_wordcloud_chart,
            create_combo_chart, create_stacked_bar_chart
        )
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        title = data.get('title', '')
        subtitle = data.get('subtitle', '')
        description = data.get('description', '')
        
        # ===== 左侧：图表区域 =====
        # 有配图时：图表占55%（严格限制在0.8in-7.3in范围内）
        # 无配图时：图表占100%
        if has_side_image:
            chart_left = Layout.MARGIN_LEFT
            chart_top = Inches(1.4)
            chart_width = Inches(6.2)  # 严格限制在7.0in以内
            chart_height = Inches(5.5)
            # 右侧配图从7.5in开始
            right_content_left = Inches(7.5)
            right_content_width = Inches(5.3)
        else:
            chart_left = Layout.MARGIN_LEFT
            chart_top = Inches(1.4)
            chart_width = Inches(7.5)
            chart_height = Inches(5.5)
            right_content_left = Inches(8.5)
            right_content_width = Inches(4.5)
        
        # 生成图表（使用白色主题）
        try:
            chart_img_path = f'/tmp/v3_{chart_type}.png'
            
            if chart_type == 'radar':
                items = data.get('items', [])
                values = [float(item.get('value', 0)) for item in items]
                labels = [f"D{i+1}" for i in range(len(items))]
                img = create_radar_chart(values, labels, "", theme=chart_theme)
                img.save(chart_img_path)
                
            elif chart_type == 'funnel':
                levels = data.get('levels', [])
                values = [float(l.get('value', 0)) for l in levels]
                labels = [l.get('title', '') for l in levels]
                img = create_funnel_chart(values, labels, "", theme=chart_theme)
                img.save(chart_img_path)
                
            elif chart_type == 'matrix':
                quads = data.get('quadrants', [])
                img = create_matrix_chart(quads, "", "", "", theme=chart_theme)
                img.save(chart_img_path)
                
            elif chart_type == 'pyramid':
                levels = data.get('levels', [])
                values = [int(l.get('title', '30').split('%')[0].strip()) if '%' in l.get('title', '') else 30 for l in levels]
                labels = [l.get('title', '') for l in levels]
                img = create_pyramid_chart(values, labels, "", theme=chart_theme)
                img.save(chart_img_path)
                
            elif chart_type == 'gantt':
                tasks = data.get('tasks', [])
                img = create_gantt_chart(tasks, "", theme=chart_theme)
                img.save(chart_img_path)
                
            elif chart_type == 'heatmap':
                matrix = data.get('matrix', [])
                xlabels = data.get('xlabels', [])
                ylabels = data.get('ylabels', [])
                if matrix and len(matrix) > 0:
                    import numpy as np
                    data_arr = np.array(matrix)
                    img = create_heatmap_chart(data_arr, xlabels, ylabels, "", theme=chart_theme)
                    img.save(chart_img_path)
                
            elif chart_type == 'gauge':
                value = float(data.get('value', 0))
                label = data.get('label', '')
                min_val = float(data.get('min', 0))
                max_val = float(data.get('max', 100))
                img = create_gauge_chart(value, label, min_val, max_val, "", theme=chart_theme)
                img.save(chart_img_path)
                
            elif chart_type == 'wordcloud':
                words = data.get('words', {})
                img = create_wordcloud_chart(words, "", theme=chart_theme)
                img.save(chart_img_path)
                
            elif chart_type == 'combo':
                categories = data.get('categories', [])
                bar_values = data.get('bar_values', [])
                line_values = data.get('line_values', [])
                img = create_combo_chart(categories, bar_values, line_values, "")
                img.save(chart_img_path)
                
            elif chart_type == 'stacked':
                categories = data.get('categories', [])
                segments = data.get('segments', {})
                img = create_stacked_bar_chart(categories, segments, "")
                img.save(chart_img_path)
            
            # 插入图表
            slide.shapes.add_picture(chart_img_path, chart_left, chart_top, width=chart_width)
            
        except Exception as e:
            print(f"图表生成失败: {e}")
            import traceback
            traceback.print_exc()
        
        # ===== 页面大标题（严格限制宽度）=====
        title_width = chart_width if has_side_image else Layout.CONTENT_WIDTH
        if title:
            title_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
                title_width, Inches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = t['text']
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(1.0),
                title_width, Inches(0.4)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = t['muted']
        
        # 右侧洞察区域（有配图时移到最右侧）
        if has_side_image:
            insight_left = right_content_left
            insight_width = right_content_width
        else:
            insight_left = right_content_left
            insight_width = right_content_width
        
        # 右侧标题
        right_title_box = slide.shapes.add_textbox(
            insight_left, chart_top,
            insight_width, Inches(0.5)
        )
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "📊 核心洞察"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        
        # 右侧分隔线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            insight_left, chart_top + Inches(0.55),
            insight_width, Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = t['border']
        line.line.fill.background()
        
        # 右侧内容 - 根据图表类型生成不同说明
        content_y = chart_top + Inches(0.7)
        
        if chart_type == 'radar':
            self._render_radar_insights(slide, data, insight_left, content_y, insight_width)
        elif chart_type == 'funnel':
            self._render_funnel_insights(slide, data, insight_left, content_y, insight_width)
        elif chart_type == 'gantt':
            self._render_gantt_insights(slide, data, insight_left, content_y, insight_width)
        elif chart_type == 'heatmap':
            self._render_heatmap_insights(slide, data, insight_left, content_y, insight_width)
        elif chart_type == 'matrix':
            self._render_matrix_insights(slide, data, insight_left, content_y, insight_width)
        elif chart_type == 'pyramid':
            self._render_pyramid_insights(slide, data, insight_left, content_y, insight_width)
        elif chart_type == 'combo':
            self._render_combo_insights(slide, data, insight_left, content_y, insight_width)
        elif chart_type == 'stacked':
            self._render_stacked_insights(slide, data, insight_left, content_y, insight_width)
        elif chart_type == 'gauge':
            self._render_gauge_insights(slide, data, insight_left, content_y, insight_width)
        else:
            # 默认渲染items
            self._render_items_legend_v2(slide, data, right_left, content_y, right_width)
        
        return slide
    
    def _render_radar_insights(self, slide, data, left, top, width):
        """雷达图 - 详细洞察"""
        t = self.t
        items = data.get('items', [])
        
        # 标题
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "能力评估分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 计算洞察
        values = [float(item.get('value', 0)) for item in items]
        avg = sum(values) / len(values) if values else 0
        max_val = max(values) if values else 0
        min_val = min(values) if values else 0
        max_item = max(items, key=lambda x: float(x.get('value', 0))) if items else {}
        min_item = min(items, key=lambda x: float(x.get('value', 0))) if items else {}
        
        y_offset = top + Inches(0.5)
        
        # 平均值
        avg_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
        tf = avg_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• 综合得分：{avg:.0f}分"
        p.font.size = Pt(14)
        p.font.color.rgb = t['text']
        y_offset += Inches(0.4)
        
        # 最高项
        if max_item:
            max_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = max_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• 优势能力：{max_item.get('title', '')}"
            p.font.size = Pt(14)
            p.font.color.rgb = t['success']
            p2 = tf.add_paragraph()
            p2.text = f"  {max_item.get('value', 0)}分，领先平均水平{max_val - avg:.0f}分"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
        
        # 最低项
        if min_item:
            min_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = min_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• 待提升：{min_item.get('title', '')}"
            p.font.size = Pt(14)
            p.font.color.rgb = t['warning']
            p2 = tf.add_paragraph()
            p2.text = f"  {min_item.get('value', 0)}分，需提升{min_val + 10 - float(min_item.get('value', 0)):.0f}分达平均"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
        
        # 建议
        suggest_box = slide.shapes.add_textbox(left, y_offset + Inches(0.3), width, Inches(0.6))
        tf = suggest_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "💡 建议"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        p2 = tf.add_paragraph()
        p2.text = f"聚焦提升{min_item.get('title', '')}能力，预计可提升综合得分{avg + 5:.0f}分"
        p2.font.size = Pt(12)
        p2.font.color.rgb = t['muted']
    
    def _render_funnel_insights(self, slide, data, left, top, width):
        """漏斗图 - 详细洞察"""
        t = self.t
        levels = data.get('levels', [])
        
        # 标题
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "转化漏斗分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        # 计算转化率
        if len(levels) >= 2:
            first_val = float(levels[0].get('value', 1))
            last_val = float(levels[-1].get('value', 0))
            total_conv = last_val / first_val * 100 if first_val > 0 else 0
            
            conv_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = conv_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• 总体转化率：{total_conv:.2f}%"
            p.font.size = Pt(14)
            p.font.color.rgb = t['text']
            y_offset += Inches(0.4)
            
            # 各阶段转化
            for i in range(len(levels) - 1):
                curr = float(levels[i].get('value', 0))
                next_val = float(levels[i+1].get('value', 0))
                conv = next_val / curr * 100 if curr > 0 else 0
                
                stage_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.35))
                tf = stage_box.text_frame
                p = tf.paragraphs[0]
                curr_name = levels[i].get('title', '').split()[0]
                next_name = levels[i+1].get('title', '').split()[0]
                p.text = f"• {curr_name}→{next_name}"
                p.font.size = Pt(13)
                p.font.color.rgb = t['muted']
                p2 = tf.add_paragraph()
                p2.text = f"  转化率 {conv:.1f}%，流失 {curr - next_val:,.0f}人"
                p2.font.size = Pt(12)
                p2.font.color.rgb = t['muted']
                y_offset += Inches(0.5)
        
        # 关键流失点
        if len(levels) >= 2:
            losses = []
            for i in range(len(levels) - 1):
                curr = float(levels[i].get('value', 0))
                next_val = float(levels[i+1].get('value', 0))
                loss = curr - next_val
                losses.append((levels[i].get('title', ''), loss))
            
            max_loss = max(losses, key=lambda x: x[1]) if losses else ('', 0)
            
            loss_box = slide.shapes.add_textbox(left, y_offset + Inches(0.2), width, Inches(0.5))
            tf = loss_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "⚠️ 最大流失点"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = t['danger']
            p2 = tf.add_paragraph()
            p2.text = f"{max_loss[0]}环节流失最多，建议重点优化"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
    
    def _render_gantt_insights(self, slide, data, left, top, width):
        """甘特图 - 详细洞察"""
        t = self.t
        tasks = data.get('tasks', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "项目进度分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        # 统计
        total = len(tasks)
        complete = len([t for t in tasks if t.get('status') == 'complete'])
        in_progress = len([t for t in tasks if t.get('status') == 'in_progress'])
        pending = len([t for t in tasks if t.get('status') == 'pending'])
        
        stats = f"• 总任务：{total}个"
        stats_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.35))
        tf = stats_box.text_frame
        p = tf.paragraphs[0]
        p.text = stats
        p.font.size = Pt(13)
        p.font.color.rgb = t['text']
        y_offset += Inches(0.35)
        
        status_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.9))
        tf = status_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"✅ 已完成：{complete}个"
        p.font.size = Pt(13)
        p.font.color.rgb = t['success']
        p2 = tf.add_paragraph()
        p2.text = f"🔄 进行中：{in_progress}个"
        p2.font.size = Pt(13)
        p2.font.color.rgb = t['accent']
        p3 = tf.add_paragraph()
        p3.text = f"⏳ 待开始：{pending}个"
        p3.font.size = Pt(13)
        p3.font.color.rgb = t['muted']
        y_offset += Inches(0.9)
        
        # 关键里程碑
        in_progress_tasks = [t for t in tasks if t.get('status') == 'in_progress']
        if in_progress_tasks:
            milestone_box = slide.shapes.add_textbox(left, y_offset + Inches(0.2), width, Inches(0.5))
            tf = milestone_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "🎯 当前重点"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = t['accent']
            for task in in_progress_tasks[:2]:
                p2 = tf.add_paragraph()
                p2.text = f"• {task.get('name', '')}（{task.get('progress', 0)}%）"
                p2.font.size = Pt(12)
                p2.font.color.rgb = t['muted']
    
    def _render_heatmap_insights(self, slide, data, left, top, width):
        """热力图 - 详细洞察"""
        t = self.t
        matrix = data.get('matrix', [])
        xlabels = data.get('xlabels', [])
        ylabels = data.get('ylabels', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "区域表现分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        if matrix and len(matrix) > 0 and len(matrix[0]) > 0:
            import numpy as np
            data_arr = np.array(matrix)
            
            # 最高值
            max_idx = np.unravel_index(data_arr.argmax(), data_arr.shape)
            max_val = data_arr.max()
            max_row = ylabels[max_idx[0]] if max_idx[0] < len(ylabels) else ''
            max_col = xlabels[max_idx[1]] if max_idx[1] < len(xlabels) else ''
            
            max_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = max_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"🏆 最佳表现"
            p.font.size = Pt(14)
            p.font.color.rgb = t['success']
            p2 = tf.add_paragraph()
            p2.text = f"{max_row} × {max_col} = {max_val:.0f}"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
            
            # 最低值
            min_idx = np.unravel_index(data_arr.argmin(), data_arr.shape)
            min_val = data_arr.min()
            min_row = ylabels[min_idx[0]] if min_idx[0] < len(ylabels) else ''
            min_col = xlabels[min_idx[1]] if min_idx[1] < len(xlabels) else ''
            
            min_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = min_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"📍 待提升区域"
            p.font.size = Pt(14)
            p.font.color.rgb = t['warning']
            p2 = tf.add_paragraph()
            p2.text = f"{min_row} × {min_col} = {min_val:.0f}，需重点关注"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
    
    def _render_matrix_insights(self, slide, data, left, top, width):
        """矩阵图 - 详细洞察"""
        t = self.t
        quads = data.get('quadrants', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "产品组合分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        # 各象限分析
        quadrant_names = ["明星产品", "问题产品", "瘦狗产品", "金牛产品"]
        colors = [t['success'], t['warning'], t['danger'], t['accent']]
        
        for i, (quadrant, name, color) in enumerate(zip(quads, quadrant_names, colors)):
            items = quadrant if isinstance(quadrant, list) else (quadrant.get('items', []) if isinstance(quadrant, dict) else [])
            if not items:
                continue
                
            q_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.35))
            tf = q_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {name}：{len(items)}个"
            p.font.size = Pt(13)
            p.font.color.rgb = color
            y_offset += Inches(0.35)
            
            if i == 0:  # 明星产品
                p2 = tf.add_paragraph()
                p2.text = "  高增长高份额，重点投入资源"
                p2.font.size = Pt(11)
                p2.font.color.rgb = t['muted']
            elif i == 3:  # 金牛产品
                p2 = tf.add_paragraph()
                p2.text = "  稳定盈利，贡献现金流"
                p2.font.size = Pt(11)
                p2.font.color.rgb = t['muted']
            y_offset += Inches(0.15)
    
    def _render_pyramid_insights(self, slide, data, left, top, width):
        """金字塔 - 详细洞察"""
        t = self.t
        levels = data.get('levels', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "品类结构分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        # 解读
        for level in levels[:3]:
            level_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.6))
            tf = level_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            title = level.get('title', '')
            desc = level.get('desc', '')
            p.text = f"• {title}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = t['accent']
            if desc:
                p2 = tf.add_paragraph()
                p2.text = f"  {desc}"
                p2.font.size = Pt(11)
                p2.font.color.rgb = t['muted']
            y_offset += Inches(0.6)
    
    def _render_combo_insights(self, slide, data, left, top, width):
        """组合图 - 详细洞察"""
        t = self.t
        bar_values = data.get('bar_values', [])
        line_values = data.get('line_values', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "趋势分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        if bar_values and line_values:
            import numpy as np
            bar_arr = np.array(bar_values)
            line_arr = np.array(line_values)
            
            # 趋势
            trend = "上升" if line_arr[-1] > line_arr[0] else "下降"
            change = line_arr[-1] - line_arr[0]
            
            trend_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = trend_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• 增长率趋势：{trend}"
            p.font.size = Pt(13)
            p.font.color.rgb = t['text']
            p2 = tf.add_paragraph()
            p2.text = f"  {change:+.1f}个百分点"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
            
            # 峰值
            max_bar_idx = np.argmax(bar_arr)
            max_bar = bar_arr.max()
            
            peak_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = peak_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• 峰值：{max_bar:.0f}"
            p.font.size = Pt(13)
            p.font.color.rgb = t['success']
            p2 = tf.add_paragraph()
            categories = data.get('categories', [])
            if max_bar_idx < len(categories):
                p2.text = f"  出现在{categories[max_bar_idx]}"
                p2.font.size = Pt(12)
                p2.font.color.rgb = t['muted']
    
    def _render_stacked_insights(self, slide, data, left, top, width):
        """堆叠图 - 详细洞察"""
        t = self.t
        segments = data.get('segments', {})
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "构成变化分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        seg_names = list(segments.keys())
        if len(seg_names) >= 2:
            first_vals = [segments[seg][0] if len(segments[seg]) > 0 else 0 for seg in seg_names]
            last_vals = [segments[seg][-1] if len(segments[seg]) > 0 else 0 for seg in seg_names]
            
            changes = [last - first for first, last in zip(first_vals, last_vals)]
            
            # 增长最快的
            max_change = max(changes)
            max_idx = changes.index(max_change)
            
            change_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.5))
            tf = change_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"📈 增长最快：{seg_names[max_idx]}"
            p.font.size = Pt(13)
            p.font.color.rgb = t['success']
            p2 = tf.add_paragraph()
            p2.text = f"  占比从{first_vals[max_idx]}%增至{last_vals[max_idx]}%（+{max_change}%）"
            p2.font.size = Pt(11)
            p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
            
            # 下降的
            min_change = min(changes)
            if min_change < 0:
                min_idx = changes.index(min_change)
                decline_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.5))
                tf = decline_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"📉 占比下降：{seg_names[min_idx]}"
                p.font.size = Pt(13)
                p.font.color.rgb = t['danger']
                p2 = tf.add_paragraph()
                p2.text = f"  占比从{first_vals[min_idx]}%降至{last_vals[min_idx]}%（{min_change}%）"
                p2.font.size = Pt(11)
                p2.font.color.rgb = t['muted']
    
    def _render_gauge_insights(self, slide, data, left, top, width):
        """仪表盘 - 详细洞察"""
        t = self.t
        value = float(data.get('value', 0))
        label = data.get('label', '')
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "KPI达成分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        # 状态评估
        if value >= 80:
            status = "优秀"
            color = t['success']
        elif value >= 60:
            status = "良好"
            color = t['accent']
        elif value >= 40:
            status = "一般"
            color = t['warning']
        else:
            status = "需改进"
            color = t['danger']
        
        status_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
        tf = status_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• 状态：{status}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        y_offset += Inches(0.4)
        
        # 解读
        gap = 100 - value
        gap_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.5))
        tf = gap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• 距目标差距"
        p.font.size = Pt(13)
        p.font.color.rgb = t['text']
        p2 = tf.add_paragraph()
        p2.text = f"还差{gap:.0f}个百分点达到100%目标"
        p2.font.size = Pt(12)
        p2.font.color.rgb = t['muted']
    
    def _render_items_legend_v2(self, slide, data, left, top, width):
        """通用items渲染"""
        t = self.t
        items = data.get('items', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "数据要点"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        for item in items[:5]:
            item_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.5))
            tf = item_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"• {item.get('title', '')}"
            p.font.size = Pt(13)
            p.font.color.rgb = t['text']
            if item.get('value'):
                p2 = tf.add_paragraph()
                p2.text = f"  {item['value']}"
                p2.font.size = Pt(12)
                p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
    
    def _render_items_legend(self, slide, items, max_val=100):
        """渲染指标图例"""
        t = self.t
        legend_x = Inches(8.8)
        legend_y = Inches(2.0)
        
        for i, item in enumerate(items):
            # 标签
            label_box = slide.shapes.add_textbox(
                legend_x, legend_y + Inches(i * 0.8),
                Inches(4), Inches(0.7)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {item.get('title', '')}"
            p.font.size = Pt(18)
            p.font.color.rgb = t['text']
            
            # 值
            val_box = slide.shapes.add_textbox(
                legend_x + Inches(0.3), legend_y + Inches(i * 0.8 + 0.35),
                Inches(3.5), Inches(0.4)
            )
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(item.get('value', ''))
            p.font.size = Pt(14)
            p.font.color.rgb = t['muted']
            
            # 进度条
            val = float(item.get('value', 0)) / max_val
            bar_width = val * 3.5
            
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                legend_x, legend_y + Inches(i * 0.8 + 0.55),
                Inches(bar_width), Inches(0.08)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = t['accent']
            bar.line.fill.background()
    
    def _render_levels_legend(self, slide, levels):
        """渲染层级图例"""
        t = self.t
        legend_x = Inches(9)
        legend_y = Inches(2.0)
        
        for i, level in enumerate(levels):
            label_box = slide.shapes.add_textbox(
                legend_x, legend_y + Inches(i * 1.0),
                Inches(4), Inches(0.9)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = level.get('title', '')
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = t['accent']
            
            if level.get('desc'):
                p2 = tf.add_paragraph()
                p2.text = level['desc']
                p2.font.size = Pt(14)
                p2.font.color.rgb = t['muted']
    
    def page_content_focus(self, title, items, style='cards', content_width_ratio=1.0, bottom_text=None):
        """📝 内容页 - 聚焦要点
        
        Args:
            content_width_ratio: 内容区域宽度比例，默认1.0（占满）
            如果右侧有配图，设为0.65（占65%宽度，右边图片35%）
            bottom_text: 底部补充文字，用于说明核心内容和补充内容
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 计算内容区域宽度（严格限制，不能超出）
        # 内容区最大宽度 = 页面宽度 - 左边距 - 右侧配图区(45%) - 间隙
        # 13.33in - 0.8in - 5.28in - 0.2in = 7.05in
        max_content_width = Inches(7.0) if content_width_ratio < 1.0 else Layout.CONTENT_WIDTH
        content_width = min(Layout.CONTENT_WIDTH * content_width_ratio, max_content_width)
        
        # 页面标题（严格限制宽度）
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
            content_width, Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        if style == 'cards':
            # 卡片式布局（严格按内容区宽度计算）
            card_width = (content_width - Inches(0.6)) / 2
            card_height = Inches(1.8)
            
            for i, item in enumerate(items):
                row = i // 2
                col = i % 2
                
                left = Layout.MARGIN_LEFT + col * (card_width + Inches(0.2))
                top = Inches(2.0) + row * (card_height + Inches(0.3))
                
                # 确保卡片不超出内容区
                if left + card_width > Layout.MARGIN_LEFT + content_width:
                    card_width = content_width - (left - Layout.MARGIN_LEFT) - Inches(0.2)
                
                # 卡片背景
                card = self._card(slide, left, top, card_width, card_height)
                
                # 标题
                item_title = item.get('title', '') if isinstance(item, dict) else item
                title_tb = slide.shapes.add_textbox(
                    left + Inches(0.3), top + Inches(0.3),
                    card_width - Inches(0.6), Inches(0.6)
                )
                tf = title_tb.text_frame
                p = tf.paragraphs[0]
                p.text = item_title
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = t['accent']
                
                # 值
                if isinstance(item, dict) and item.get('value'):
                    val_tb = slide.shapes.add_textbox(
                        left + Inches(0.3), top + Inches(0.9),
                        card_width - Inches(0.6), Inches(0.5)
                    )
                    tf = val_tb.text_frame
                    p = tf.paragraphs[0]
                    p.text = item['value']
                    p.font.size = Pt(18)
                    p.font.color.rgb = t['text']
        
        # 底部补充文字
        if bottom_text:
            bottom_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(6.5),
                Layout.CONTENT_WIDTH, Inches(0.8)
            )
            tf = bottom_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = bottom_text
            p.font.size = Pt(14)
            p.font.color.rgb = t['muted']
            p.font.italic = True
        
        return slide
    
    def page_icons_grid(self, title, items, columns=3, sub=None, description=None):
        """📋 图标网格页 - 用于目录或概览
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
            Layout.CONTENT_WIDTH, Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 副标题
        if sub:
            sub_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Layout.MARGIN_TOP + Inches(0.8),
                Layout.CONTENT_WIDTH, Inches(0.5)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = sub
            p.font.size = Pt(18)
            p.font.color.rgb = t['muted']
        
        # 描述
        desc_start_y = Layout.MARGIN_TOP + Inches(1.3)
        if description:
            desc_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, desc_start_y,
                Layout.CONTENT_WIDTH, Inches(0.6)
            )
            tf = desc_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = description
            p.font.size = Pt(14)
            p.font.color.rgb = t['muted']
            desc_start_y += Inches(0.7)
        
        # 图标网格
        num_items = len(items)
        card_width = (Layout.CONTENT_WIDTH - Inches(0.3) * (columns - 1)) / columns
        card_height = Inches(1.5)
        start_y = desc_start_y + Inches(0.3)
        start_x = Layout.MARGIN_LEFT
        
        for i, item in enumerate(items):
            row = i // columns
            col = i % columns
            
            x = start_x + col * (card_width + Inches(0.3))
            y = start_y + row * (card_height + Inches(0.2))
            
            # 卡片背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = t['card_bg']
            card.line.fill.background()
            card.adjustments[0] = 0.05
            
            # 图标和标题
            icon = item.get('icon', '📌')
            item_title = item.get('title', '')
            item_desc = item.get('desc', '')
            
            # 图标
            icon_box = slide.shapes.add_textbox(
                x + Inches(0.15), y + Inches(0.15),
                Inches(0.6), Inches(0.6)
            )
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = icon
            p.font.size = Pt(28)
            
            # 标题
            title_text = slide.shapes.add_textbox(
                x + Inches(0.15), y + Inches(0.7),
                card_width - Inches(0.3), Inches(0.4)
            )
            tf = title_text.text_frame
            p = tf.paragraphs[0]
            p.text = item_title
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = t['text']
            
            # 描述
            if item_desc:
                desc_text = slide.shapes.add_textbox(
                    x + Inches(0.15), y + Inches(1.05),
                    card_width - Inches(0.3), Inches(0.35)
                )
                tf = desc_text.text_frame
                p = tf.paragraphs[0]
                p.text = item_desc
                p.font.size = Pt(12)
                p.font.color.rgb = t['muted']
        
        return slide

    def page_quote(self, quote, attribution=None):
        """💬 引用页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 大引号
        quote_mark = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Inches(1.5),
            Inches(2), Inches(2)
        )
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = '"'
        p.font.size = Pt(200)
        p.font.color.rgb = t['accent']
        p.font.name = 'Georgia'
        
        # 引用内容
        quote_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.8),
            Layout.CONTENT_WIDTH - Inches(1), Inches(3)
        )
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = quote
        p.font.size = Pt(32)
        p.font.italic = True
        p.font.color.rgb = t['text']
        
        # 作者
        if attribution:
            attr_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(6),
                Layout.CONTENT_WIDTH, Inches(0.6)
            )
            tf = attr_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {attribution}"
            p.font.size = Pt(20)
            p.font.color.rgb = t['muted']
            p.alignment = PP_ALIGN.RIGHT
        
        return slide
    
    def page_timeline(self, title, steps):
        """📅 时间线页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
            Layout.CONTENT_WIDTH, Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 时间线
        num = len(steps)
        start_x = Layout.MARGIN_LEFT + Inches(1)
        end_x = Layout.WIDTH - Layout.MARGIN_RIGHT - Inches(1)
        line_y = Inches(4.2)
        
        # 主线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            start_x, line_y, end_x - start_x, Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = t['border']
        line.line.fill.background()
        
        # 节点和内容
        for i, step in enumerate(steps):
            x = start_x + (end_x - start_x) * i / max(num - 1, 1)
            
            # 节点圆
            node = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x - Inches(0.2), line_y - Inches(0.2),
                Inches(0.4), Inches(0.4)
            )
            node.fill.solid()
            node.fill.fore_color.rgb = t['accent']
            node.line.fill.background()
            
            # 步骤标题
            step_title = step.get('title', f'Step {i+1}')
            st_box = slide.shapes.add_textbox(
                x - Inches(1), line_y + Inches(0.5),
                Inches(2), Inches(0.6)
            )
            tf = st_box.text_frame
            p = tf.paragraphs[0]
            p.text = step_title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = t['text']
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            if step.get('desc'):
                desc_box = slide.shapes.add_textbox(
                    x - Inches(1.2), line_y + Inches(1.1),
                    Inches(2.4), Inches(1)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = step['desc']
                p.font.size = Pt(14)
                p.font.color.rgb = t['muted']
                p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def page_process_flow(self, title, steps, sub=None):
        """🔄 流程图页 - 用于展示流程步骤
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
            Layout.CONTENT_WIDTH, Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 副标题
        if sub:
            sub_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Layout.MARGIN_TOP + Inches(0.7),
                Layout.CONTENT_WIDTH, Inches(0.4)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = sub
            p.font.size = Pt(16)
            p.font.color.rgb = t['muted']
        
        # 流程步骤
        num = len(steps)
        start_x = Layout.MARGIN_LEFT + Inches(0.5)
        end_x = Layout.WIDTH - Layout.MARGIN_RIGHT - Inches(0.5)
        box_width = (end_x - start_x - Inches(0.3) * (num - 1)) / num
        box_height = Inches(1.8)
        start_y = Inches(2.8)
        
        for i, step in enumerate(steps):
            x = start_x + (box_width + Inches(0.3)) * i
            
            # 步骤框
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, start_y, box_width, box_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = t['card_bg']
            box.line.color.rgb = t['border']
            
            # 步骤序号
            num_circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x + box_width/2 - Inches(0.25), start_y - Inches(0.25),
                Inches(0.5), Inches(0.5)
            )
            num_circle.fill.solid()
            num_circle.fill.fore_color.rgb = t['accent']
            num_circle.line.fill.background()
            
            # 序号文字
            num_box = slide.shapes.add_textbox(
                x + box_width/2 - Inches(0.25), start_y - Inches(0.22),
                Inches(0.5), Inches(0.5)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = t['text']
            p.alignment = PP_ALIGN.CENTER
            
            # 步骤标题
            step_title = step.get('title', f'Step {i+1}')
            title_box = slide.shapes.add_textbox(
                x + Inches(0.1), start_y + Inches(0.3),
                box_width - Inches(0.2), Inches(0.6)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = step_title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = t['text']
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            if step.get('desc'):
                desc_box = slide.shapes.add_textbox(
                    x + Inches(0.1), start_y + Inches(0.9),
                    box_width - Inches(0.2), Inches(0.8)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = step['desc']
                p.font.size = Pt(11)
                p.font.color.rgb = t['muted']
                p.alignment = PP_ALIGN.CENTER
            
            # 箭头（除了最后一个）
            if i < num - 1:
                arrow_x = x + box_width + Inches(0.05)
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    arrow_x, start_y + box_height/2 - Inches(0.15),
                    Inches(0.2), Inches(0.3)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = t['accent']
                arrow.line.fill.background()
        
        return slide
    
    def page_icons_grid(self, title, items, columns=4, sub=None):
        """🔲 图标网格页 - 带图标的要点网格展示
        
        Args:
            title: 页面标题
            items: 要点列表 [{"title": "标题", "desc": "描述", "icon": "emoji"}, ...]
            columns: 列数（默认4列）
            sub: 副标题
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
            Layout.CONTENT_WIDTH, Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 副标题
        if sub:
            sub_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Layout.MARGIN_TOP + Inches(0.7),
                Layout.CONTENT_WIDTH, Inches(0.4)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = sub
            p.font.size = Pt(16)
            p.font.color.rgb = t['muted']
        
        # 计算网格布局
        num = len(items)
        rows = (num + columns - 1) // columns
        
        margin_left = Layout.MARGIN_LEFT + Inches(0.3)
        margin_top = Inches(1.8)
        box_width = (Layout.CONTENT_WIDTH - Inches(0.6) * (columns - 1)) / columns
        box_height = Inches(1.5)
        gap_x = Inches(0.6)
        gap_y = Inches(0.3)
        
        for i, item in enumerate(items):
            row = i // columns
            col = i % columns
            
            x = margin_left + col * (box_width + gap_x)
            y = margin_top + row * (box_height + gap_y)
            
            # 卡片背景
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, box_width, box_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = t['card_bg']
            card.line.color.rgb = t['border']
            
            # 图标（如果提供了emoji）
            icon = item.get('icon', '●')
            icon_box = slide.shapes.add_textbox(
                x, y + Inches(0.15),
                box_width, Inches(0.5)
            )
            tf = icon_box.text_frame
            p = tf.paragraphs[0]
            p.text = icon
            p.font.size = Pt(28)
            p.alignment = PP_ALIGN.CENTER
            
            # 标题
            item_title = item.get('title', '')
            title_box = slide.shapes.add_textbox(
                x + Inches(0.1), y + Inches(0.6),
                box_width - Inches(0.2), Inches(0.4)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item_title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = t['text']
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            if item.get('desc'):
                desc_box = slide.shapes.add_textbox(
                    x + Inches(0.1), y + Inches(1.0),
                    box_width - Inches(0.2), Inches(0.4)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item['desc']
                p.font.size = Pt(11)
                p.font.color.rgb = t['muted']
                p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def page_comparison(self, title, left_items, right_items, left_title="方案A", right_title="方案B", sub=None):
        """⚖️ 对比图页 - 左右两栏对比展示
        
        Args:
            title: 页面标题
            left_items: 左侧要点 [{"title": "标题", "value": "值"}, ...]
            right_items: 右侧要点
            left_title: 左侧标题
            right_title: 右侧标题
            sub: 副标题
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
            Layout.CONTENT_WIDTH, Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 副标题
        if sub:
            sub_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Layout.MARGIN_TOP + Inches(0.7),
                Layout.CONTENT_WIDTH, Inches(0.4)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = sub
            p.font.size = Pt(16)
            p.font.color.rgb = t['muted']
        
        # 左右栏布局
        col_width = (Layout.CONTENT_WIDTH - Inches(0.6)) / 2
        left_x = Layout.MARGIN_LEFT
        right_x = Layout.MARGIN_LEFT + col_width + Inches(0.6)
        start_y = Inches(1.8)
        card_height = Inches(0.7)
        gap_y = Inches(0.15)
        
        # 左侧标题
        left_title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_x, start_y, col_width, Inches(0.5)
        )
        left_title_box.fill.solid()
        left_title_box.fill.fore_color.rgb = t['accent']
        left_title_box.line.fill.background()
        
        lt_box = slide.shapes.add_textbox(
            left_x, start_y + Inches(0.05),
            col_width, Inches(0.4)
        )
        tf = lt_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 右侧标题
        right_title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            right_x, start_y, col_width, Inches(0.5)
        )
        right_title_box.fill.solid()
        right_title_box.fill.fore_color.rgb = t['warning']
        right_title_box.line.fill.background()
        
        rt_box = slide.shapes.add_textbox(
            right_x, start_y + Inches(0.05),
            col_width, Inches(0.4)
        )
        tf = rt_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 左侧内容
        for i, item in enumerate(left_items):
            y = start_y + Inches(0.65) + i * (card_height + gap_y)
            
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left_x, y, col_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = t['card_bg']
            card.line.color.rgb = t['border']
            
            # 标题
            t_box = slide.shapes.add_textbox(
                left_x + Inches(0.15), y + Inches(0.1),
                col_width - Inches(0.3), Inches(0.3)
            )
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = item.get('title', '')
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = t['text']
            
            # 值
            v_box = slide.shapes.add_textbox(
                left_x + Inches(0.15), y + Inches(0.35),
                col_width - Inches(0.3), Inches(0.3)
            )
            tf = v_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(item.get('value', ''))
            p.font.size = Pt(12)
            p.font.color.rgb = t['muted']
        
        # 右侧内容
        for i, item in enumerate(right_items):
            y = start_y + Inches(0.65) + i * (card_height + gap_y)
            
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                right_x, y, col_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = t['card_bg']
            card.line.color.rgb = t['border']
            
            # 标题
            t_box = slide.shapes.add_textbox(
                right_x + Inches(0.15), y + Inches(0.1),
                col_width - Inches(0.3), Inches(0.3)
            )
            tf = t_box.text_frame
            p = tf.paragraphs[0]
            p.text = item.get('title', '')
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = t['text']
            
            # 值
            v_box = slide.shapes.add_textbox(
                right_x + Inches(0.15), y + Inches(0.35),
                col_width - Inches(0.3), Inches(0.3)
            )
            tf = v_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(item.get('value', ''))
            p.font.size = Pt(12)
            p.font.color.rgb = t['muted']
        
        return slide
    
    def page_section(self, title, subtitle=None):
        """📑 章节分隔页 - 简洁的章节标题页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 左侧强调线
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Layout.MARGIN_LEFT, Inches(2.8),
            Inches(0.12), Inches(1.8)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = t['accent']
        accent.line.fill.background()
        
        # 主标题 - 黑色
        title_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(2.8),
            Layout.CONTENT_WIDTH - Inches(0.4), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = t['text']  # V4.9.1修复：使用主题色
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(4.3),
                Layout.CONTENT_WIDTH, Inches(0.8)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = t['muted']
        
        return slide

    def page_closing(self, title, cta=None, contact=None):
        """🎬 结尾页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 主标题 - 居中大标题
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Inches(2.5),
            Layout.CONTENT_WIDTH, Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = t['text']
        p.alignment = PP_ALIGN.CENTER
        
        # CTA
        if cta:
            cta_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(4.5),
                Layout.CONTENT_WIDTH, Inches(0.8)
            )
            tf = cta_box.text_frame
            p = tf.paragraphs[0]
            p.text = cta
            p.font.size = Pt(24)
            p.font.color.rgb = t['accent']
            p.alignment = PP_ALIGN.CENTER
        
        # 联系
        if contact:
            contact_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(6),
                Layout.CONTENT_WIDTH, Inches(0.5)
            )
            tf = contact_box.text_frame
            p = tf.paragraphs[0]
            p.text = contact
            p.font.size = Pt(16)
            p.font.color.rgb = t['muted']
            p.alignment = PP_ALIGN.CENTER
        
        return slide
def create_ppt_v3(slides, output_path, theme='executive'):
    """V3.0创建PPT"""
    gen = PPTGeneratorV3(theme)
    
    for slide in slides:
        stype = slide.get('type', '')
        
        if stype == 'title':
            gen.page_title(
                slide.get('main', ''),
                slide.get('sub'),
                slide.get('tag')
            )
        elif stype == 'big_number':
            gen.page_big_number(
                slide.get('number', ''),
                slide.get('label', ''),
                slide.get('context')
            )
        elif stype == 'two_columns':
            gen.page_two_columns(
                slide.get('left_title', ''),
                slide.get('left_items', []),
                slide.get('right_title', ''),
                slide.get('right_items', []),
                slide.get('page_title')
            )
        elif stype == 'radar':
            gen.page_chart_visual(slide, 'radar')
        elif stype == 'funnel':
            gen.page_chart_visual(slide, 'funnel')
        elif stype == 'matrix':
            gen.page_chart_visual(slide, 'matrix')
        elif stype == 'pyramid':
            gen.page_chart_visual(slide, 'pyramid')
        elif stype == 'gantt':
            gen.page_chart_visual(slide, 'gantt')
        elif stype == 'heatmap':
            gen.page_chart_visual(slide, 'heatmap')
        elif stype == 'gauge':
            gen.page_chart_visual(slide, 'gauge')
        elif stype == 'wordcloud':
            gen.page_chart_visual(slide, 'wordcloud')
        elif stype == 'combo':
            gen.page_chart_visual(slide, 'combo')
        elif stype == 'stacked':
            gen.page_chart_visual(slide, 'stacked')
        elif stype == 'section':
            gen.page_section(
                slide.get('title', ''),
                slide.get('subtitle')
            )
        elif stype == 'icons_grid':
            gen.page_icons_grid(
                slide.get('title', ''),
                slide.get('items', []),
                slide.get('columns', 3),
                slide.get('sub')
            )
        elif stype == 'content_cards':
            gen.page_content_focus(
                slide.get('title', ''),
                slide.get('items', []),
                'cards'
            )
        elif stype == 'quote':
            gen.page_quote(
                slide.get('quote', ''),
                slide.get('attribution')
            )
        elif stype == 'timeline':
            gen.page_timeline(
                slide.get('title', ''),
                slide.get('steps', [])
            )
        elif stype == 'closing':
            gen.page_closing(
                slide.get('title', ''),
                slide.get('cta'),
                slide.get('contact')
            )
    
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    gen.prs.save(output_path)
    return output_path
def create_ppt_v3_with_images(slides, output_path, theme='executive', page_images=None, page_image_configs=None, auto_plan_images=True):
    """V3.0创建PPT - 支持AI配图融合版本 V3.0
    
    Args:
        slides: 幻灯片配置列表
        output_path: 输出路径
        theme: 配色主题
        page_images: dict, {page_index: image_path} 简单的图片路径映射
        page_image_configs: dict, {page_index: {
            "path": str,           # 图片路径
            "position": str,       # 位置 (background/left/right/top_right/bottom_right/header/wide_banner)
            "size": str,          # 尺寸 (full/large/medium/small/tiny)
            "opacity": float,     # 透明度
            "add_overlay": bool   # 是否添加遮罩
        }} 详细的配图配置
        auto_plan_images: bool, 是否自动规划配图（使用image_planner）
    """
    page_images = page_images or {}
    page_image_configs = page_image_configs or {}
    
    # 导入配图规划模块
    if auto_plan_images:
        try:
            from image_planner import plan_page_images
            topic_hint = slides[0].get('main', '') if slides and isinstance(slides[0], dict) else ''
            auto_strategies = plan_page_images(slides, topic_hint)
            
            # 合并配置：用户配置优先，自动规划作为补充
            for idx, strategy in auto_strategies.items():
                if idx not in page_image_configs:
                    # 使用自动规划的策略（但用户没有配置图片路径）
                    # 这里只记录策略，实际图片生成由外部调用
                    pass
        except ImportError:
            pass  # 配图规划模块不可用
    
    gen = PPTGeneratorV3(theme)
    
    for idx, slide in enumerate(slides):
        stype = slide.get('type', '')
        
        # 获取配图配置
        img_config = page_image_configs.get(idx, {})
        image_path = img_config.get("path") or page_images.get(idx)
        
        if stype == 'title':
            # 封面：全屏背景
            gen.page_title_with_image(
                slide.get('main', ''),
                slide.get('sub'),
                slide.get('tag'),
                image_path
            )
        elif stype == 'big_number':
            # 大数据页：右侧装饰
            gen.page_big_number_with_image(
                slide.get('number', ''),
                slide.get('label', ''),
                slide.get('context'),
                image_path
            )
        elif stype == 'closing':
            # 结尾：全屏背景
            gen.page_closing_with_image(
                slide.get('title', ''),
                slide.get('cta'),
                slide.get('contact'),
                image_path
            )
        elif stype == 'section':
            # 章节分隔页：图片作为右下角装饰，不覆盖内容
            gen.page_section(
                slide.get('title', ''),
                slide.get('subtitle')
            )
            if image_path and os.path.exists(image_path):
                gen._add_image_ex(gen.prs.slides[idx], image_path, "decorative_right", "medium")
        elif stype == 'icons_grid':
            # 图标网格页（如目录页默认不显示配图）
            gen.page_icons_grid(
                slide.get('title', ''),
                slide.get('items', []),
                slide.get('columns', 3),
                slide.get('sub')
            )
            # 只有明确设置 show_image=True 时才显示配图
            if image_path and os.path.exists(image_path) and slide.get('show_image'):
                gen._add_image_ex(gen.prs.slides[idx], image_path, "decorative_right", "medium")
        elif stype == 'two_columns':
            # 双栏：右侧装饰
            gen.page_two_columns(
                slide.get('left_title', ''),
                slide.get('left_items', []),
                slide.get('right_title', ''),
                slide.get('right_items', []),
                slide.get('page_title')
            )
            if image_path and os.path.exists(image_path):
                position = img_config.get("position", "right")
                size = img_config.get("size", "medium")
                gen._add_image_ex(gen.prs.slides[idx], image_path, position, size)
        elif stype == 'content_cards':
            # 卡片页：有配图时内容占65%（左侧），无配图时占100%
            has_image = image_path and os.path.exists(image_path)
            width_ratio = 0.55 if has_image else 1.0
            gen.page_content_focus(
                slide.get('title', ''),
                slide.get('items', []),
                'cards',
                content_width_ratio=width_ratio,
                bottom_text=slide.get('bottom_text')
            )
            if has_image:
                # 根据 size 选择合适的位置
                img_size = img_config.get("size", "small")
                if img_size == "small":
                    position = "right_small"  # 30%宽度，靠右
                elif img_size == "tiny":
                    position = "right_tiny"  # 25%宽度
                else:
                    position = "right_medium"  # 45%宽度
                gen._add_image_ex(gen.prs.slides[idx], image_path, position, img_size)
        elif stype == 'radar':
            # 雷达图：有配图时图表缩小
            has_image = image_path and os.path.exists(image_path)
            gen.page_chart_visual(slide, 'radar', has_side_image=has_image)
            if has_image:
                # 根据 size 选择合适的位置
                img_size = img_config.get("size", "small")
                if img_size == "small":
                    position = "right_small"  # 30%宽度，靠右
                elif img_size == "tiny":
                    position = "right_tiny"  # 25%宽度
                else:
                    position = "right_medium"  # 45%宽度
                gen._add_image_ex(gen.prs.slides[idx], image_path, position, img_size)
        elif stype == 'funnel':
            has_image = image_path and os.path.exists(image_path)
            gen.page_chart_visual(slide, 'funnel', has_side_image=has_image)
            if has_image:
                # 根据 size 选择合适的位置
                img_size = img_config.get("size", "small")
                if img_size == "small":
                    position = "right_small"  # 30%宽度，靠右
                elif img_size == "tiny":
                    position = "right_tiny"  # 25%宽度
                else:
                    position = "right_medium"  # 45%宽度
                gen._add_image_ex(gen.prs.slides[idx], image_path, position, img_size)
        elif stype == 'matrix':
            gen.page_chart_visual(slide, 'matrix')
        elif stype == 'pyramid':
            gen.page_chart_visual(slide, 'pyramid')
        elif stype == 'gantt':
            gen.page_chart_visual(slide, 'gantt')
        elif stype == 'heatmap':
            gen.page_chart_visual(slide, 'heatmap')
        elif stype == 'gauge':
            has_image = image_path and os.path.exists(image_path)
            gen.page_chart_visual(slide, 'gauge', has_side_image=has_image)
            if has_image:
                # 根据 size 选择合适的位置
                img_size = img_config.get("size", "small")
                if img_size == "small":
                    position = "right_small"  # 30%宽度，靠右
                elif img_size == "tiny":
                    position = "right_tiny"  # 25%宽度
                else:
                    position = "right_medium"  # 45%宽度
                gen._add_image_ex(gen.prs.slides[idx], image_path, position, img_size)
        elif stype == 'wordcloud':
            gen.page_chart_visual(slide, 'wordcloud')
        elif stype == 'combo':
            gen.page_chart_visual(slide, 'combo')
        elif stype == 'stacked':
            gen.page_chart_visual(slide, 'stacked')
        elif stype == 'quote':
            gen.page_quote(
                slide.get('quote', ''),
                slide.get('attribution')
            )
            if image_path and os.path.exists(image_path):
                gen._add_image_ex(gen.prs.slides[idx], image_path, "left", "medium")
        elif stype == 'timeline':
            gen.page_timeline(
                slide.get('title', ''),
                slide.get('steps', [])
            )
            if image_path and os.path.exists(image_path):
                gen._add_image_ex(gen.prs.slides[idx], image_path, "right", "small")
    
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    gen.prs.save(output_path)
    return output_path
