#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT生成器 - 经典经典版 V1.0
专为传统老板/领导设计的PPT风格
配色：深红+金色 | 排版：对称工整 | 图表：经典柱状图/饼图/折线图
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu
import io
from PIL import Image, ImageDraw, ImageFont
import numpy as np

# ═══════════════════════════════════════════════════════════════════════════════
# 经典经典配色主题
# ═══════════════════════════════════════════════════════════════════════════════

def to_rgb(color_tuple):
    """RGB元组转matplotlib颜色"""
    if isinstance(color_tuple, tuple) and len(color_tuple) == 3:
        return '#{:02x}{:02x}{:02x}'.format(*[int(c) for c in color_tuple])
    return color_tuple

THEMES = {
    "classic": {
        "bg": (45, 27, 27),           # 深棕红背景
        "card_bg": (60, 40, 40),      # 卡片背景
        "text": (255, 248, 230),       # 米白色文字
        "muted": (180, 170, 160),      # 浅灰文字
        "accent": (218, 165, 32),      # 金色强调（经典）
        "light_accent": (255, 215, 0), # 亮金色
        "success": (107, 142, 35),     # 橄榄绿
        "warning": (184, 134, 11),     # 深金色
        "danger": (178, 34, 34),      # 深红色
        "border": (100, 80, 60),      # 边框色
        "title_bg": (60, 30, 30),     # 标题栏背景
    },
}

DEFAULT_THEME = "classic"

# ═══════════════════════════════════════════════════════════════════════════════
# 经典图表生成器（经典风格）
# ═══════════════════════════════════════════════════════════════════════════════

def create_bar_chart(data, labels, title="", theme=DEFAULT_THEME):
    """经典柱状图 - 经典最爱"""
    import matplotlib.pyplot as plt
    
    t = THEMES.get(theme, THEMES[DEFAULT_THEME])
    
    fig, ax = plt.subplots(figsize=(10, 6), facecolor=to_rgb(t['bg']))
    ax.set_facecolor(to_rgb(t['bg']))
    
    x = np.arange(len(labels))
    bars = ax.bar(x, data, color=to_rgb(t['accent']), alpha=0.85, width=0.6)
    
    # 数值标注
    for bar, val in zip(bars, data):
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{val:,.0f}',
                ha='center', va='bottom', fontsize=12, 
                fontweight='bold', color=to_rgb(t['text']))
    
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=12, color=to_rgb(t['text']))
    ax.tick_params(colors=to_rgb(t['muted']))
    
    # 网格线（淡淡的）
    ax.yaxis.grid(True, color=to_rgb(t['border']), alpha=0.3, linewidth=0.5)
    ax.set_axisbelow(True)
    
    # 隐藏边框
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.spines['bottom'].set_color(to_rgb(t['border']))
    ax.spines['left'].set_color(to_rgb(t['border']))
    
    if title:
        ax.set_title(title, fontsize=18, fontweight='bold', color=to_rgb(t['text']), pad=15)
    
    plt.tight_layout()
    return fig_to_image(fig)


def create_pie_chart(sizes, labels, title="", theme=DEFAULT_THEME):
    """经典饼图 - 构成分析"""
    import matplotlib.pyplot as plt
    
    t = THEMES.get(theme, THEMES[DEFAULT_THEME])
    
    fig, ax = plt.subplots(figsize=(8, 8), facecolor=to_rgb(t['bg']))
    ax.set_facecolor(to_rgb(t['bg']))
    
    colors = [to_rgb(t['accent']), to_rgb(t['success']), to_rgb(t['warning']), to_rgb(t['danger']), to_rgb(t['border'])]
    
    # 经典饼图：实心、无旋转、显示百分比
    wedges, texts, autotexts = ax.pie(
        sizes, 
        labels=labels,
        colors=colors[:len(sizes)],
        autopct='%1.1f%%',
        startangle=90,
        textprops={'color': to_rgb(t['text']), 'fontsize': 12},
        wedgeprops={'edgecolor': to_rgb(t['bg']), 'linewidth': 2}
    )
    
    # 百分比文字样式
    for autotext in autotexts:
        autotext.set_color(to_rgb(t['bg']))
        autotext.set_fontweight('bold')
        autotext.set_fontsize(14)
    
    if title:
        ax.set_title(title, fontsize=18, fontweight='bold', color=to_rgb(t['text']), pad=20)
    
    # 保持圆形
    ax.axis('equal')
    plt.tight_layout()
    return fig_to_image(fig)


def create_line_chart(data, labels, title="", theme=DEFAULT_THEME):
    """经典折线图 - 趋势分析"""
    import matplotlib.pyplot as plt
    
    t = THEMES.get(theme, THEMES[DEFAULT_THEME])
    
    fig, ax = plt.subplots(figsize=(10, 6), facecolor=to_rgb(t['bg']))
    ax.set_facecolor(to_rgb(t['bg']))
    
    x = np.arange(len(labels))
    
    # 经典折线：粗线 + 实心点 + 标注
    ax.plot(x, data, 
            color=to_rgb(t['accent']), 
            linewidth=3,
            marker='o', 
            markersize=12,
            markerfacecolor=to_rgb(t['light_accent']),
            markeredgecolor=to_rgb(t['accent']),
            markeredgewidth=2)
    
    # 数值标注
    for i, val in enumerate(data):
        ax.annotate(f'{val:,.0f}', 
                   (i, val), 
                   textcoords="offset points", 
                   xytext=(0, 15), 
                   ha='center', 
                   fontsize=11,
                   fontweight='bold',
                   color=to_rgb(t['text']))
    
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=12, color=to_rgb(t['text']))
    ax.tick_params(colors=to_rgb(t['muted']))
    
    # 网格
    ax.yaxis.grid(True, color=to_rgb(t['border']), alpha=0.3, linewidth=0.5)
    ax.set_axisbelow(True)
    
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.spines['bottom'].set_color(to_rgb(t['border']))
    ax.spines['left'].set_color(to_rgb(t['border']))
    
    if title:
        ax.set_title(title, fontsize=18, fontweight='bold', color=to_rgb(t['text']), pad=15)
    
    plt.tight_layout()
    return fig_to_image(fig)


def create_comparison_chart(data1, data2, labels, title="", theme=DEFAULT_THEME):
    """经典对比柱状图 - 同比/环比"""
    import matplotlib.pyplot as plt
    
    t = THEMES.get(theme, THEMES[DEFAULT_THEME])
    
    fig, ax = plt.subplots(figsize=(10, 6), facecolor=to_rgb(t['bg']))
    ax.set_facecolor(to_rgb(t['bg']))
    
    x = np.arange(len(labels))
    width = 0.35
    
    bars1 = ax.bar(x - width/2, data1, width, label=f'2025年', color=to_rgb(t['accent']), alpha=0.85)
    bars2 = ax.bar(x + width/2, data2, width, label=f'2024年', color=to_rgb(t['muted']), alpha=0.7)
    
    # 数值标注
    for bar in bars1:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{height:,.0f}',
                ha='center', va='bottom', fontsize=10, color=to_rgb(t['text']))
    
    for bar in bars2:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{height:,.0f}',
                ha='center', va='bottom', fontsize=10, color=to_rgb(t['muted']))
    
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=12, color=to_rgb(t['text']))
    ax.tick_params(colors=to_rgb(t['muted']))
    ax.legend(fontsize=12, framealpha=0.8, facecolor=to_rgb(t['card_bg']), edgecolor=to_rgb(t['border']))
    
    ax.yaxis.grid(True, color=to_rgb(t['border']), alpha=0.3, linewidth=0.5)
    ax.set_axisbelow(True)
    
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.spines['bottom'].set_color(to_rgb(t['border']))
    ax.spines['left'].set_color(to_rgb(t['border']))
    
    if title:
        ax.set_title(title, fontsize=18, fontweight='bold', color=to_rgb(t['text']), pad=15)
    
    plt.tight_layout()
    return fig_to_image(fig)


def create_kpi_gauge(value, label="", max_val=100, title="", theme=DEFAULT_THEME):
    """经典KPI仪表盘 - 达成率"""
    import matplotlib.pyplot as plt
    
    t = THEMES.get(theme, THEMES[DEFAULT_THEME])
    
    fig, ax = plt.subplots(figsize=(6, 4), facecolor=to_rgb(t['bg']))
    ax.set_facecolor(to_rgb(t['bg']))
    
    # 简单进度条样式
    ratio = min(value / max_val, 1.0) if max_val > 0 else 0
    
    # 背景条
    ax.barh(0.5, 1, height=0.3, left=0, color=to_rgb(t['border']), alpha=0.3)
    
    # 进度条
    color = to_rgb(t['success']) if ratio >= 0.8 else (to_rgb(t['warning']) if ratio >= 0.5 else to_rgb(t['danger']))
    ax.barh(0.5, ratio, height=0.3, left=0, color=color, alpha=0.9)
    
    # 中心文字
    ax.text(0.5, 0.5, f'{value:.0f}%', fontsize=36, fontweight='bold',
           ha='center', va='center', color=to_rgb(t['text']))
    ax.text(0.5, 0.1, label, fontsize=14, ha='center', va='center', color=to_rgb(t['muted']))
    
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis('off')
    
    if title:
        ax.set_title(title, fontsize=16, fontweight='bold', color=to_rgb(t['text']), pad=10)
    
    plt.tight_layout()
    return fig_to_image(fig)


def create_table_image(data, title="", theme=DEFAULT_THEME):
    """经典数据表格"""
    import matplotlib.pyplot as plt
    
    t = THEMES.get(theme, THEMES[DEFAULT_THEME])
    
    rows = len(data) + 1  # +1 header
    cols = len(data[0]) if data else 3
    
    fig, ax = plt.subplots(figsize=(10, rows * 0.6 + 1))
    ax.set_facecolor(to_rgb(t['bg']))
    fig.patch.set_facecolor(to_rgb(t['bg']))
    
    ax.axis('off')
    
    # 表头背景
    header_rect = plt.Rectangle((0, 1 - 1/rows), 1, 1/rows, 
                                facecolor=to_rgb(t['title_bg']), edgecolor=to_rgb(t['border']))
    ax.add_patch(header_rect)
    
    # 单元格
    for i, row in enumerate(data):
        y_pos = 1 - (i + 1.5) / rows
        
        for j, cell in enumerate(row):
            x_pos = j / cols
            
            # 单元格背景
            cell_rect = plt.Rectangle((x_pos, y_pos - 0.5/rows), 
                                     1/cols, 1/rows,
                                     facecolor=to_rgb(t['card_bg']) if i % 2 == 0 else to_rgb(t['bg']),
                                     edgecolor=to_rgb(t['border']),
                                     linewidth=0.5)
            ax.add_patch(cell_rect)
            
            # 文字
            if i == 0:  # 表头
                ax.text(x_pos + 0.5/cols, y_pos, str(cell),
                       ha='center', va='center',
                       fontsize=11, fontweight='bold', color=to_rgb(t['text']))
            else:  # 数据
                ax.text(x_pos + 0.5/cols, y_pos, str(cell),
                       ha='center', va='center',
                       fontsize=10, color=to_rgb(t['text']))
    
    if title:
        ax.text(0.5, 1.05, title, transform=ax.transAxes,
               fontsize=14, fontweight='bold', 
               ha='center', va='bottom', color=to_rgb(t['text']))
    
    plt.tight_layout()
    return fig_to_image(fig)


def fig_to_image(fig):
    """matplotlib转PIL Image"""
    import matplotlib.pyplot as plt
    import io
    from PIL import Image
    
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
               facecolor=fig.get_facecolor())
    buf.seek(0)
    img = Image.open(buf).convert('RGB')
    plt.close(fig)
    return img


# ═══════════════════════════════════════════════════════════════════════════════
# PPT页面生成器
# ═══════════════════════════════════════════════════════════════════════════════

class ClassicPPTGer:
    """经典经典风格PPT生成器"""
    
    def __init__(self, theme='classic'):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.t = THEMES.get(theme, THEMES[DEFAULT_THEME])
    
    def _color(self, key):
        """获取RGBColor对象"""
        c = self.t.get(key, (255, 255, 255))
        if isinstance(c, RGBColor):
            return c
        return RGBColor(c[0], c[1], c[2])
    
    def _add_bg(self, slide):
        """添加背景"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._color('bg')
    
    def _add_header_bar(self, slide, title):
        """经典风格：顶部金色标题栏"""
        # 标题栏背景
        header = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._color('title_bg')
        header.line.fill.background()
        
        # 金色下划线
        line = slide.shapes.add_shape(
            1,
            Inches(0), Inches(1.15),
            self.prs.slide_width, Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._color('accent')
        line.line.fill.background()
        
        # 标题文字
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.3),
            Inches(11), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._color('accent')
    
    def _add_footer(self, slide, page_num, total):
        """经典风格：底部页码"""
        footer_box = slide.shapes.add_textbox(
            Inches(12), Inches(7),
            Inches(1), Inches(0.3)
        )
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{page_num}/{total}"
        p.font.size = Pt(12)
        p.font.color.rgb = self._color('muted')
        p.alignment = PP_ALIGN.RIGHT
    
    def page_title(self, data):
        """封面页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        
        # 中央大标题（对称布局）
        main_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.33), Inches(1.5)
        )
        tf = main_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get('main', '')
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self._color('text')
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        sub_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.2),
            Inches(11.33), Inches(0.6)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get('sub', '')
        p.font.size = Pt(24)
        p.font.color.rgb = self._color('muted')
        p.alignment = PP_ALIGN.CENTER
        
        # 底部信息
        tag_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.5),
            Inches(11.33), Inches(0.5)
        )
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get('tag', '')
        p.font.size = Pt(16)
        p.font.color.rgb = self._color('accent')
        p.alignment = PP_ALIGN.CENTER
        
        # 装饰线
        line = slide.shapes.add_shape(
            1,
            Inches(5), Inches(5.2),
            Inches(3.33), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._color('accent')
        line.line.fill.background()
        
        return slide
    
    def page_big_number(self, data):
        """大数据页 - 对称布局"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_header_bar(slide, data.get('title', '核心成果'))
        
        # 中心大数字
        num_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5),
            Inches(11.33), Inches(2)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get('number', '0')
        p.font.size = Pt(120)
        p.font.bold = True
        p.font.color.rgb = self._color('accent')
        p.alignment = PP_ALIGN.CENTER
        
        # 标签
        label_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5),
            Inches(11.33), Inches(0.6)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get('label', '')
        p.font.size = Pt(28)
        p.font.color.rgb = self._color('text')
        p.alignment = PP_ALIGN.CENTER
        
        # 上下文
        ctx_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.3),
            Inches(11.33), Inches(0.5)
        )
        tf = ctx_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get('context', '')
        p.font.size = Pt(18)
        p.font.color.rgb = self._color('muted')
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def page_bar_chart(self, data):
        """柱状图页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_header_bar(slide, data.get('title', ''))
        
        try:
            categories = data.get('categories', [])
            values = data.get('values', [])
            
            if categories and values:
                img = create_bar_chart(values, categories, "", "classic")
                img_path = '/tmp/classic_bar.png'
                img.save(img_path)
                slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.8), width=Inches(10))
                
                # 右侧说明
                insights = data.get('insights', [])
                if insights:
                    y_pos = 5.5
                    for insight in insights[:3]:
                        box = slide.shapes.add_textbox(
                            Inches(10), Inches(y_pos),
                            Inches(2.5), Inches(0.4)
                        )
                        tf = box.text_frame
                        p = tf.paragraphs[0]
                        p.text = f"◆ {insight}"
                        p.font.size = Pt(11)
                        p.font.color.rgb = self._color('muted')
                        y_pos += 0.4
                        
        except Exception as e:
            print(f"柱状图生成失败: {e}")
        
        return slide
    
    def page_pie_chart(self, data):
        """饼图页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_header_bar(slide, data.get('title', ''))
        
        try:
            segments = data.get('segments', {})
            if segments:
                labels = list(segments.keys())
                sizes = list(segments.values())
                
                img = create_pie_chart(sizes, labels, "", "classic")
                img_path = '/tmp/classic_pie.png'
                img.save(img_path)
                slide.shapes.add_picture(img_path, Inches(2), Inches(1.5), width=Inches(8))
                
        except Exception as e:
            print(f"饼图生成失败: {e}")
        
        return slide
    
    def page_line_chart(self, data):
        """折线图页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_header_bar(slide, data.get('title', ''))
        
        try:
            categories = data.get('categories', [])
            values = data.get('values', [])
            
            if categories and values:
                img = create_line_chart(values, categories, "", "classic")
                img_path = '/tmp/classic_line.png'
                img.save(img_path)
                slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.8), width=Inches(10))
                
        except Exception as e:
            print(f"折线图生成失败: {e}")
        
        return slide
    
    def page_comparison(self, data):
        """对比图页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_header_bar(slide, data.get('title', ''))
        
        try:
            categories = data.get('categories', [])
            data_2025 = data.get('data_2025', [])
            data_2024 = data.get('data_2024', [])
            
            if categories and data_2025 and data_2024:
                img = create_comparison_chart(data_2025, data_2024, categories, "", "classic")
                img_path = '/tmp/classic_compare.png'
                img.save(img_path)
                slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.8), width=Inches(10))
                
        except Exception as e:
            print(f"对比图生成失败: {e}")
        
        return slide
    
    def page_kpi_dashboard(self, data):
        """KPI仪表盘页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_header_bar(slide, data.get('title', '核心KPI达成情况'))
        
        items = data.get('items', [])
        n = len(items)
        
        # 2x2 布局
        positions = [(2, 2), (7, 2), (2, 4.5), (7, 4.5)]
        
        for i, item in enumerate(items[:4]):
            if i >= 4:
                break
                
            x, y = positions[i]
            
            # KPI卡片
            card = slide.shapes.add_shape(
                1,
                Inches(x), Inches(y),
                Inches(4.5), Inches(2)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self._color('card_bg')
            card.line.color.rgb = self._color('border')
            
            # 数值
            value_text = str(item.get('value', '0%'))
            val_box = slide.shapes.add_textbox(
                Inches(x + 0.3), Inches(y + 0.3),
                Inches(3.9), Inches(1)
            )
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = value_text
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = self._color('accent')
            
            # 标签
            title_text = item.get('title', '')
            label_box = slide.shapes.add_textbox(
                Inches(x + 0.3), Inches(y + 1.3),
                Inches(3.9), Inches(0.5)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(14)
            p.font.color.rgb = self._color('text')
        
        return slide
    
    def page_table(self, data):
        """数据表格页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_header_bar(slide, data.get('title', ''))
        
        try:
            headers = data.get('headers', [])
            rows = data.get('rows', [])
            
            if headers or rows:
                table_data = [headers] + rows
                img = create_table_image(table_data, "", "classic")
                img_path = '/tmp/classic_table.png'
                img.save(img_path)
                slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.8), width=Inches(10))
                
        except Exception as e:
            print(f"表格生成失败: {e}")
        
        return slide
    
    def page_content_list(self, data):
        """经典风格内容列表"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_header_bar(slide, data.get('title', ''))
        
        items = data.get('items', [])
        
        # 左对齐内容
        y_pos = 1.8
        for item in items[:6]:
            # 圆点
            dot = slide.shapes.add_shape(
                9,  # OVAL
                Inches(1.5), Inches(y_pos + 0.15),
                Inches(0.15), Inches(0.15)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = self._color('accent')
            dot.line.fill.background()
            
            # 文字
            text_box = slide.shapes.add_textbox(
                Inches(1.8), Inches(y_pos),
                Inches(9), Inches(0.6)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item.get('title', '')
            p.font.size = Pt(20)
            p.font.color.rgb = self._color('text')
            
            # 值
            if item.get('value'):
                p2 = tf.add_paragraph()
                p2.text = item['value']
                p2.font.size = Pt(16)
                p2.font.color.rgb = self._color('accent')
            
            y_pos += 0.8
        
        return slide
    
    def page_two_columns(self, data):
        """经典风格双栏对比"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        
        # 顶部标题
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.3),
            Inches(11.73), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get('page_title', '')
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._color('accent')
        
        # 分隔线
        line = slide.shapes.add_shape(
            1,
            Inches(0.8), Inches(1.1),
            Inches(11.73), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._color('accent')
        line.line.fill.background()
        
        # 左栏
        left_title = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5),
            Inches(5.5), Inches(0.5)
        )
        tf = left_title.text_frame
        p = tf.paragraphs[0]
        p.text = data.get('left_title', '')
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self._color('danger')
        
        left_items = data.get('left_items', [])
        y_pos = 2.2
        for item in left_items[:5]:
            dot = slide.shapes.add_shape(
                9,
                Inches(1), Inches(y_pos + 0.15),
                Inches(0.12), Inches(0.12)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = self._color('danger')
            dot.line.fill.background()
            
            text_box = slide.shapes.add_textbox(
                Inches(1.3), Inches(y_pos),
                Inches(5), Inches(0.5)
            )
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = self._color('text')
            y_pos += 0.7
        
        # 右栏
        right_title = slide.shapes.add_textbox(
            Inches(7), Inches(1.5),
            Inches(5.5), Inches(0.5)
        )
        tf = right_title.text_frame
        p = tf.paragraphs[0]
        p.text = data.get('right_title', '')
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = self._color('success')
        
        right_items = data.get('right_items', [])
        y_pos = 2.2
        for item in right_items[:5]:
            dot = slide.shapes.add_shape(
                9,
                Inches(7.2), Inches(y_pos + 0.15),
                Inches(0.12), Inches(0.12)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = self._color('success')
            dot.line.fill.background()
            
            text_box = slide.shapes.add_textbox(
                Inches(7.5), Inches(y_pos),
                Inches(5), Inches(0.5)
            )
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = self._color('text')
            y_pos += 0.7
        
        return slide
    
    def page_quote(self, data):
        """经典风格金句页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        
        # 中央引号装饰
        quote_mark = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(1), Inches(1)
        )
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = '"'
        p.font.size = Pt(120)
        p.font.color.rgb = self._color('accent')
        
        # 引用文字
        quote_box = slide.shapes.add_textbox(
            Inches(2), Inches(3),
            Inches(9.33), Inches(2)
        )
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get('quote', '')
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.color.rgb = self._color('text')
        p.alignment = PP_ALIGN.CENTER
        
        # 来源
        attr_box = slide.shapes.add_textbox(
            Inches(2), Inches(5.2),
            Inches(9.33), Inches(0.5)
        )
        tf = attr_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"— {data.get('attribution', '')}"
        p.font.size = Pt(18)
        p.font.color.rgb = self._color('muted')
        p.alignment = PP_ALIGN.RIGHT
        
        return slide
    
    def page_closing(self, data):
        """结尾页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        
        # 中央标题
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.8),
            Inches(11.33), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get('title', '谢谢观看')
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self._color('text')
        p.alignment = PP_ALIGN.CENTER
        
        # 装饰线
        line = slide.shapes.add_shape(
            1,
            Inches(5), Inches(4),
            Inches(3.33), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._color('accent')
        line.line.fill.background()
        
        # 副标题
        cta_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.3),
            Inches(11.33), Inches(0.6)
        )
        tf = cta_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get('cta', '')
        p.font.size = Pt(20)
        p.font.color.rgb = self._color('muted')
        p.alignment = PP_ALIGN.CENTER
        
        # 部门
        contact_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.2),
            Inches(11.33), Inches(0.5)
        )
        tf = contact_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get('contact', '')
        p.font.size = Pt(16)
        p.font.color.rgb = self._color('accent')
        p.alignment = PP_ALIGN.CENTER
        
        return slide


def create_classic_ppt(slides, output_path, theme='classic'):
    """创建经典经典风格PPT"""
    ger = ClassicPPTGer(theme)
    
    for i, slide_data in enumerate(slides, 1):
        slide_type = slide_data.get('type', 'content')
        
        if slide_type == 'title':
            ger.page_title(slide_data)
        elif slide_type == 'big_number':
            ger.page_big_number(slide_data)
        elif slide_type == 'bar':
            ger.page_bar_chart(slide_data)
        elif slide_type == 'pie':
            ger.page_pie_chart(slide_data)
        elif slide_type == 'line':
            ger.page_line_chart(slide_data)
        elif slide_type == 'comparison':
            ger.page_comparison(slide_data)
        elif slide_type == 'kpi':
            ger.page_kpi_dashboard(slide_data)
        elif slide_type == 'table':
            ger.page_table(slide_data)
        elif slide_type == 'content_list':
            ger.page_content_list(slide_data)
        elif slide_type == 'two_columns':
            ger.page_two_columns(slide_data)
        elif slide_type == 'quote':
            ger.page_quote(slide_data)
        elif slide_type == 'closing':
            ger.page_closing(slide_data)
        else:
            # 默认内容页
            ger.page_content_list(slide_data)
    
    ger.prs.save(output_path)
    return output_path


if __name__ == '__main__':
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    
    # 测试
    slides = [
        {'type': 'title', 'main': '2025年度工作汇报', 'sub': '稳中求进 再创佳绩', 'tag': '市场部 · 2026年1月'},
        {'type': 'big_number', 'number': '1280万', 'label': '年度销售额', 'context': '同比增长 34.7%'},
        {'type': 'bar', 'title': '季度销售额趋势', 'categories': ['Q1', 'Q2', 'Q3', 'Q4'], 'values': [280, 320, 350, 330]},
        {'type': 'closing', 'title': '谢谢观看', 'cta': '不忘初心 砥砺前行', 'contact': '市场部'}
    ]
    
    output = create_classic_ppt(slides, '/tmp/classic_test.pptx')
    print(f"经典版PPT已生成: {output}")
