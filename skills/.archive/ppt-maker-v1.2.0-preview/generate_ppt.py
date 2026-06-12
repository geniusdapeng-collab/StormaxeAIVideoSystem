#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT生成器 - V2.0（专业级）
支持4种专业主题 + 11种幻灯片模板
配色参考 AI Presentation Maker
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import sys
import os

# 添加图表生成器路径
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

# ══════════════════════════════════════════════════════════════
# 专业配色主题
# ══════════════════════════════════════════════════════════════

THEMES = {
    "terminal": {
        "name": "Terminal",
        "bg": RGBColor(26, 26, 26),
        "card_bg": RGBColor(37, 37, 37),
        "text": RGBColor(255, 255, 255),
        "muted": RGBColor(179, 179, 179),
        "accent": RGBColor(0, 230, 118),
        "accent_dark": RGBColor(0, 200, 83),
        "border": RGBColor(51, 51, 51),
    },
    "executive": {
        "name": "Executive",
        "bg": RGBColor(13, 27, 42),
        "card_bg": RGBColor(27, 40, 56),
        "text": RGBColor(255, 255, 255),
        "muted": RGBColor(136, 153, 170),
        "accent": RGBColor(255, 183, 0),
        "accent_dark": RGBColor(229, 165, 0),
        "border": RGBColor(45, 63, 82),
    },
    "spark": {
        "name": "Spark",
        "bg": RGBColor(15, 12, 41),
        "card_bg": RGBColor(30, 25, 60),
        "text": RGBColor(240, 240, 240),
        "muted": RGBColor(160, 176, 192),
        "accent": RGBColor(36, 198, 220),
        "accent_dark": RGBColor(81, 74, 157),
        "border": RGBColor(60, 50, 100),
    },
    "clean": {
        "name": "Clean",
        "bg": RGBColor(255, 255, 255),
        "card_bg": RGBColor(248, 248, 248),
        "text": RGBColor(26, 26, 26),
        "muted": RGBColor(102, 102, 102),
        "accent": RGBColor(230, 57, 70),
        "accent_dark": RGBColor(197, 48, 60),
        "border": RGBColor(224, 224, 224),
    },
}

DEFAULT_THEME = "spark"


def get_theme(theme_name=None):
    """获取主题配色"""
    return THEMES.get(theme_name or DEFAULT_THEME, THEMES[DEFAULT_THEME])


def rgb_to_hex(rgb):
    """RGB转HEX"""
    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


class PPTGenerator:
    """专业PPT生成器"""
    
    def __init__(self, theme=DEFAULT_THEME):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9
        self.prs.slide_height = Inches(7.5)
        self.theme = get_theme(theme)
    
    # ══════════════════════════════════════════════════════════
    # 基础布局方法
    # ══════════════════════════════════════════════════════════
    
    def _add_header_bar(self, slide, title, subtitle=None):
        """添加专业标题栏"""
        t = self.theme
        
        # 渐变效果用纯色模拟
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = t['bg']
        header.line.fill.background()
        
        # 强调线
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(1.18), self.prs.slide_width, Inches(0.04)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = t['accent']
        accent_line.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(12), Inches(0.4))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(16)
            p.font.color.rgb = t['muted']
    
    def _add_background_gradient(self, slide):
        """添加背景（深色主题）"""
        t = self.theme
        
        # 全页背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = t['bg']
        bg.line.fill.background()
        
        # 移到最底层
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
    
    # ══════════════════════════════════════════════════════════
    # 幻灯片类型
    # ══════════════════════════════════════════════════════════
    
    def add_title_slide(self, title, subtitle=None, speaker=None):
        """标题页 - 专业开场"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        # 全页背景
        self._add_background_gradient(slide)
        
        # 装饰性强调块
        accent_block = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(3.2), Inches(0.3), Inches(1.2)
        )
        accent_block.fill.solid()
        accent_block.fill.fore_color.rgb = t['accent']
        accent_block.line.fill.background()
        
        # 主标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11), Inches(0.6))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = t['muted']
        
        # 演讲者
        if speaker:
            speaker_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11), Inches(0.5))
            tf = speaker_box.text_frame
            p = tf.paragraphs[0]
            p.text = speaker
            p.font.size = Pt(18)
            p.font.color.rgb = t['accent']
        
        return slide
    
    def add_section_slide(self, title, subtitle=None):
        """章节分隔页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        
        # 大号章节标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = t['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(0.6))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(20)
            p.font.color.rgb = t['muted']
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(self, title, items, subtitle=None):
        """内容页 - 带要点列表"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        self._add_header_bar(slide, title, subtitle)
        
        # 内容区域
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            if isinstance(item, dict):
                item_title = item.get('title', '')
                item_desc = item.get('desc', '')
                item_value = item.get('value', '')
                
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"▸ {item_title}"
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = t['accent']
                p.space_before = Pt(20)
                p.space_after = Pt(5)
                
                if item_value:
                    p2 = tf.add_paragraph()
                    p2.text = f"  {item_value}"
                    p2.font.size = Pt(18)
                    p2.font.color.rgb = t['muted']
                
                if item_desc:
                    p3 = tf.add_paragraph()
                    p3.text = f"  {item_desc}"
                    p3.font.size = Pt(16)
                    p3.font.color.rgb = t['muted']
            else:
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"▸ {item}"
                p.font.size = Pt(22)
                p.font.color.rgb = t['text']
                p.space_before = Pt(15)
                p.space_after = Pt(8)
        
        return slide
    
    def add_big_number_slide(self, number, label, context=None):
        """大数据页 - 单个醒目数字"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        
        # 大号数字
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(3))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(180)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        p.alignment = PP_ALIGN.CENTER
        
        # 标签
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12), Inches(0.8))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(36)
        p.font.color.rgb = t['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 上下文
        if context:
            ctx_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(0.6))
            tf = ctx_box.text_frame
            p = tf.paragraphs[0]
            p.text = context
            p.font.size = Pt(20)
            p.font.color.rgb = t['muted']
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_comparison_slide(self, title, left_title, left_items, right_title, right_items, subtitle=None):
        """对比页 - 左右两栏"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        self._add_header_bar(slide, title, subtitle)
        
        # 左栏标题
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.5), Inches(0.5))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        
        # 左栏内容
        left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.5), Inches(4.5))
        tf = left_content.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = t['text']
            p.space_before = Pt(12)
        
        # 中间分隔线
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.6), Inches(0.02), Inches(5.5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = t['border']
        divider.line.fill.background()
        
        # 右栏标题
        right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.5), Inches(0.5))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        
        # 右栏内容
        right_content = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.5), Inches(4.5))
        tf = right_content.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = t['text']
            p.space_before = Pt(12)
        
        return slide
    
    def add_timeline_slide(self, title, steps, subtitle=None):
        """时间线页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        self._add_header_bar(slide, title, subtitle)
        
        # 时间线
        num_steps = len(steps)
        start_x = Inches(1)
        end_x = Inches(12)
        y_pos = Inches(3.8)
        
        # 主线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, start_x, y_pos, end_x - start_x, Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = t['border']
        line.line.fill.background()
        
        # 步骤
        for i, step in enumerate(steps):
            x_pos = start_x + (end_x - start_x) * i / max(num_steps - 1, 1)
            
            # 节点
            node = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, x_pos - Inches(0.15), y_pos - Inches(0.15), Inches(0.3), Inches(0.3)
            )
            node.fill.solid()
            node.fill.fore_color.rgb = t['accent']
            node.line.fill.background()
            
            # 步骤标题
            step_box = slide.shapes.add_textbox(
                x_pos - Inches(1), y_pos + Inches(0.4), Inches(2), Inches(1.5)
            )
            tf = step_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = step.get('title', f'Step {i+1}')
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = t['text']
            p.alignment = PP_ALIGN.CENTER
            
            if step.get('desc'):
                p2 = tf.add_paragraph()
                p2.text = step['desc']
                p2.font.size = Pt(12)
                p2.font.color.rgb = t['muted']
                p2.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_quote_slide(self, quote_text, attribution=None):
        """引用页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        
        # 引号装饰
        quote_mark = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2), Inches(2))
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = '"'
        p.font.size = Pt(180)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        p.font.name = 'Georgia'
        
        # 引用内容
        quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(3))
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = quote_text
        p.font.size = Pt(32)
        p.font.italic = True
        p.font.color.rgb = t['text']
        p.alignment = PP_ALIGN.CENTER
        
        # 作者
        if attribution:
            attr_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.6))
            tf = attr_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {attribution}"
            p.font.size = Pt(20)
            p.font.color.rgb = t['accent']
            p.alignment = PP_ALIGN.RIGHT
        
        return slide
    
    def add_closing_slide(self, title, cta_text=None, links=None, contact=None):
        """结尾CTA页"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        
        # 主标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = t['text']
        p.alignment = PP_ALIGN.CENTER
        
        # CTA
        if cta_text:
            cta_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(0.8))
            tf = cta_box.text_frame
            p = tf.paragraphs[0]
            p.text = cta_text
            p.font.size = Pt(28)
            p.font.color.rgb = t['accent']
            p.alignment = PP_ALIGN.CENTER
        
        # 链接
        if links:
            links_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(0.5))
            tf = links_box.text_frame
            for i, link in enumerate(links):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = link
                p.font.size = Pt(16)
                p.font.color.rgb = t['muted']
                p.alignment = PP_ALIGN.CENTER
        
        # 联系信息
        if contact:
            contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12), Inches(0.5))
            tf = contact_box.text_frame
            p = tf.paragraphs[0]
            p.text = contact
            p.font.size = Pt(16)
            p.font.color.rgb = t['muted']
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    # ══════════════════════════════════════════════════════════
    # 图表幻灯片（集成高质量图表）
    # ══════════════════════════════════════════════════════════
    
    def add_radar_slide(self, title, subtitle, items, max_value=100, theme=None):
        """雷达图幻灯片"""
        from chart_generator import create_radar_chart
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        self._add_header_bar(slide, title, subtitle)
        
        # 生成雷达图
        try:
            values = [float(item.get('value', 0)) for item in items]
            labels = [item.get('title', '') for item in items]
            en_labels = [f"D{i+1}" for i in range(len(labels))]
            
            img = create_radar_chart(values, en_labels, "", theme or DEFAULT_THEME)
            img_path = '/tmp/radar_chart.png'
            img.save(img_path)
            
            # 插入图表
            slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.8), width=Inches(7))
            
            # 右侧图例
            legend_x = Inches(8.2)
            legend_y = Inches(2.2)
            
            for i, item in enumerate(items):
                legend_box = slide.shapes.add_textbox(
                    legend_x, legend_y + Inches(i * 0.7), Inches(4.5), Inches(0.6)
                )
                tf = legend_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"• {item.get('title', '')}: {item.get('value', 0)}"
                p.font.size = Pt(16)
                p.font.color.rgb = t['text']
                
                # 进度条
                value = float(item.get('value', 0)) / max_value
                bar_width = value * 3
                
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    legend_x, legend_y + Inches(i * 0.7 + 0.35),
                    Inches(bar_width), Inches(0.15)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = t['accent']
                bar.line.fill.background()
                
        except Exception as e:
            print(f"雷达图生成失败: {e}")
            self._add_simple_content_fallback(slide, items)
        
        return slide
    
    def add_funnel_slide(self, title, subtitle, levels, theme=None):
        """漏斗图幻灯片"""
        from chart_generator import create_funnel_chart
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        self._add_header_bar(slide, title, subtitle)
        
        # 生成漏斗图
        try:
            values = [float(level.get('value', 0)) for level in levels]
            labels = [level.get('title', '') for level in levels]
            
            img = create_funnel_chart(values, labels, "", theme or DEFAULT_THEME)
            img_path = '/tmp/funnel_chart.png'
            img.save(img_path)
            
            slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.6), width=Inches(10))
            
        except Exception as e:
            print(f"漏斗图生成失败: {e}")
            self._add_simple_funnel(slide, levels)
        
        return slide
    
    def add_matrix_slide(self, title, subtitle, quadrants, theme=None):
        """矩阵图幻灯片"""
        from chart_generator import create_matrix_chart
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        self._add_header_bar(slide, title, subtitle)
        
        # 生成矩阵图
        try:
            xlabel = subtitle.split('vs')[-1].strip() if 'vs' in subtitle else 'X'
            ylabel = subtitle.split('vs')[0].strip() if 'vs' in subtitle else 'Y'
            
            quad_items = [q.get('items', []) if isinstance(q, dict) else [] for q in quadrants]
            img = create_matrix_chart(quad_items, xlabel, ylabel, "", theme or DEFAULT_THEME)
            img_path = '/tmp/matrix_chart.png'
            img.save(img_path)
            
            slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.6), width=Inches(7.5))
            
            # 右侧内容
            content_x = Inches(8.5)
            content_y = Inches(1.8)
            quadrant_names = ["明星", "问题", "瘦狗", "金牛"]
            quadrant_colors = [t['accent'], RGBColor(245, 124, 0), RGBColor(198, 40, 40), RGBColor(21, 101, 192)]
            
            for i, (quad, name) in enumerate(zip(quadrants, quadrant_names)):
                if isinstance(quad, dict) and 'items' in quad:
                    items = quad['items']
                else:
                    items = quad if isinstance(quad, list) else []
                
                q_box = slide.shapes.add_textbox(content_x, content_y + Inches(i * 1.3), Inches(4.5), Inches(1.2))
                tf = q_box.text_frame
                tf.word_wrap = True
                
                p = tf.paragraphs[0]
                p.text = f"【{name}】"
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = quadrant_colors[i]
                
                for item in items[:3]:
                    p2 = tf.add_paragraph()
                    p2.text = f"• {item}"
                    p2.font.size = Pt(14)
                    p2.font.color.rgb = t['text']
            
        except Exception as e:
            print(f"矩阵图生成失败: {e}")
            self._add_simple_matrix(slide, quadrants)
        
        return slide
    
    def add_pyramid_slide(self, title, subtitle, levels, theme=None):
        """金字塔幻灯片"""
        from chart_generator import create_pyramid_chart
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        self._add_header_bar(slide, title, subtitle)
        
        # 生成金字塔图
        try:
            # 解析百分比
            values = []
            labels = []
            for level in levels:
                title_str = level.get('title', '30%')
                # 提取数字
                import re
                match = re.search(r'(\d+)', title_str)
                val = int(match.group(1)) if match else 30
                values.append(val)
                labels.append(title_str)
            
            img = create_pyramid_chart(values, labels, "", theme or DEFAULT_THEME)
            img_path = '/tmp/pyramid_chart.png'
            img.save(img_path)
            
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.6), width=Inches(8))
            
            # 右侧说明
            if levels and any(l.get('desc') for l in levels):
                desc_x = Inches(9)
                desc_y = Inches(2)
                
                for i, level in enumerate(levels):
                    if level.get('desc'):
                        desc_box = slide.shapes.add_textbox(
                            desc_x, desc_y + Inches(i * 1.2), Inches(4), Inches(1)
                        )
                        tf = desc_box.text_frame
                        tf.word_wrap = True
                        p = tf.paragraphs[0]
                        p.text = level.get('title', '').split()[0] if ' ' in level.get('title', '') else level.get('title', '')
                        p.font.size = Pt(18)
                        p.font.bold = True
                        p.font.color.rgb = t['accent']
                        
                        p2 = tf.add_paragraph()
                        p2.text = level.get('desc', '')
                        p2.font.size = Pt(14)
                        p2.font.color.rgb = t['muted']
            
        except Exception as e:
            print(f"金字塔图生成失败: {e}")
            self._add_simple_pyramid(slide, levels)
        
        return slide
    
    def add_chart_slide(self, title, subtitle, items, theme=None):
        """图表页 - 柱状图"""
        from chart_generator import create_bar_chart
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        t = self.theme
        
        self._add_background_gradient(slide)
        self._add_header_bar(slide, title, subtitle)
        
        # 生成柱状图
        try:
            categories = [item.get('title', f'C{i+1}') for i, item in enumerate(items)]
            values = [float(item.get('value', 0)) for item in items]
            
            img = create_bar_chart(categories, values, "", theme or DEFAULT_THEME)
            img_path = '/tmp/bar_chart.png'
            img.save(img_path)
            
            slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.6), width=Inches(11.5))
            
        except Exception as e:
            print(f"图表生成失败: {e}")
            self._add_simple_content_fallback(slide, items)
        
        return slide
    
    # ══════════════════════════════════════════════════════════
    # 回退方法
    # ══════════════════════════════════════════════════════════
    
    def _add_simple_content_fallback(self, slide, items):
        """简单内容回退"""
        t = self.theme
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(11.5), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(item, dict):
                p.text = f"▸ {item.get('title', '')}: {item.get('value', '')}"
            else:
                p.text = f"▸ {item}"
            p.font.size = Pt(22)
            p.font.color.rgb = t['text']
            p.space_before = Pt(15)
    
    def _add_simple_funnel(self, slide, levels):
        """简单漏斗回退"""
        t = self.theme
        num_levels = len(levels)
        
        for i, level in enumerate(levels):
            ratio = 1 - (i * 0.15)
            width = 8 * ratio
            left = (13.333 - width) / 2
            
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left), Inches(2 + i * 0.9),
                Inches(width), Inches(0.7)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = t['accent']
            bar.line.fill.background()
            
            text_box = slide.shapes.add_textbox(
                Inches(left), Inches(2.15 + i * 0.9),
                Inches(width), Inches(0.5)
            )
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{level.get('title', '')}: {level.get('value', '')}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
    
    def _add_simple_matrix(self, slide, quadrants):
        """简单矩阵回退"""
        t = self.theme
        matrix_left = Inches(1)
        matrix_top = Inches(2.2)
        matrix_width = Inches(6)
        matrix_height = Inches(4.5)
        
        positions = [
            (matrix_left, matrix_top, "明星"),
            (matrix_left + matrix_width/2, matrix_top, "问题"),
            (matrix_left + matrix_width/2, matrix_top + matrix_height/2, "瘦狗"),
            (matrix_left, matrix_top + matrix_height/2, "金牛"),
        ]
        
        for q_left, q_top, name in positions:
            q_width = matrix_width / 2 - Inches(0.1)
            q_height = matrix_height / 2 - Inches(0.1)
            
            quad = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, q_left, q_top, q_width, q_height
            )
            quad.fill.solid()
            quad.fill.fore_color.rgb = t['card_bg']
            quad.line.color.rgb = t['border']
            
            label = slide.shapes.add_textbox(q_left + Inches(0.2), q_top + Inches(0.2), Inches(2.5), Inches(0.6))
            tf = label.text_frame
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = t['accent']
    
    def _add_simple_pyramid(self, slide, levels):
        """简单金字塔回退"""
        t = self.theme
        num_levels = len(levels)
        
        for i, level in enumerate(levels):
            ratio = 0.3 + (i * 0.7 / max(num_levels - 1, 1))
            width = 10 * ratio
            left = (13.333 - width) / 2
            
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(2.2 + i * 0.9),
                Inches(width), Inches(0.7)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = t['accent']
            bar.line.fill.background()
            
            text_box = slide.shapes.add_textbox(
                Inches(left), Inches(2.35 + i * 0.9),
                Inches(width), Inches(0.5)
            )
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = level.get('title', '')
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER


def create_ppt(slides, output_path, theme=None):
    """根据配置生成PPT"""
    generator = PPTGenerator(theme)
    
    for slide in slides:
        slide_type = slide.get('type', 'content')
        
        if slide_type == 'title':
            generator.add_title_slide(
                slide.get('title', ''),
                slide.get('subtitle'),
                slide.get('speaker')
            )
        elif slide_type == 'section':
            generator.add_section_slide(
                slide.get('title', ''),
                slide.get('subtitle')
            )
        elif slide_type == 'content':
            generator.add_content_slide(
                slide.get('title', ''),
                slide.get('items', []),
                slide.get('subtitle')
            )
        elif slide_type == 'big_number':
            generator.add_big_number_slide(
                slide.get('number', '0'),
                slide.get('label', ''),
                slide.get('context')
            )
        elif slide_type == 'comparison':
            generator.add_comparison_slide(
                slide.get('title', ''),
                slide.get('left_title', ''),
                slide.get('left_items', []),
                slide.get('right_title', ''),
                slide.get('right_items', []),
                slide.get('subtitle')
            )
        elif slide_type == 'timeline':
            generator.add_timeline_slide(
                slide.get('title', ''),
                slide.get('steps', []),
                slide.get('subtitle')
            )
        elif slide_type == 'quote':
            generator.add_quote_slide(
                slide.get('quote_text', ''),
                slide.get('attribution')
            )
        elif slide_type == 'closing':
            generator.add_closing_slide(
                slide.get('title', ''),
                slide.get('cta_text'),
                slide.get('links'),
                slide.get('contact')
            )
        elif slide_type == 'chart':
            generator.add_chart_slide(
                slide.get('title', ''),
                slide.get('subtitle', ''),
                slide.get('items', [])
            )
        elif slide_type == 'radar':
            generator.add_radar_slide(
                slide.get('title', ''),
                slide.get('subtitle', ''),
                slide.get('items', []),
                slide.get('max_value', 100)
            )
        elif slide_type == 'funnel':
            generator.add_funnel_slide(
                slide.get('title', ''),
                slide.get('subtitle', ''),
                slide.get('levels', [])
            )
        elif slide_type == 'matrix':
            generator.add_matrix_slide(
                slide.get('title', ''),
                slide.get('subtitle', ''),
                slide.get('quadrants', [])
            )
        elif slide_type == 'pyramid':
            generator.add_pyramid_slide(
                slide.get('title', ''),
                slide.get('subtitle', ''),
                slide.get('levels', [])
            )
    
    # 保存
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    generator.prs.save(output_path)
    return output_path


if __name__ == "__main__":
    # 测试
    slides = [
        {'type': 'title', 'title': '专业PPT演示', 'subtitle': '基于AI Presentation Maker设计'},
        {'type': 'section', 'title': '核心功能'},
        {'type': 'content', 'title': '四大图表类型', 'items': [
            {'title': '雷达图', 'value': '多维度能力评估'},
            {'title': '漏斗图', 'value': '转化流程可视化'},
            {'title': '矩阵图', 'value': '四象限分析'},
            {'title': '金字塔', 'value': '层级结构展示'},
        ]},
        {'type': 'radar', 'title': '员工能力评估', 'subtitle': '六维度全面分析', 'items': [
            {'title': '沟通', 'value': 85},
            {'title': '技术', 'value': 72},
            {'title': '创新', 'value': 78},
            {'title': '协作', 'value': 90},
            {'title': '执行', 'value': 88},
        ]},
        {'type': 'funnel', 'title': '用户转化漏斗', 'subtitle': '从曝光到付费', 'levels': [
            {'title': '曝光', 'value': 100000},
            {'title': '点击', 'value': 15000},
            {'title': '注册', 'value': 3000},
            {'title': '下单', 'value': 600},
            {'title': '支付', 'value': 450},
        ]},
        {'type': 'closing', 'title': '谢谢观看', 'cta_text': '立即体验PPT制作技能'},
    ]
    
    output = create_ppt(slides, '/tmp/professional_ppt.pptx')
    print(f"专业PPT已生成: {output}")
