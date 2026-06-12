#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT生成器 - V3.0（美学升级版）
核心原则：层次感 / 留白 / 图形化 / 节奏 / 网格 / 克制
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import sys
import os

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))


# ══════════════════════════════════════════════════════════════
# 排版常量 - 网格系统
# ══════════════════════════════════════════════════════════════

class Layout:
    """排版网格系统"""
    # 页面尺寸 16:9
    WIDTH = Inches(13.333)
    HEIGHT = Inches(7.5)
    
    # 边距系统
    MARGIN_LEFT = Inches(0.8)
    MARGIN_RIGHT = Inches(0.8)
    MARGIN_TOP = Inches(0.6)
    MARGIN_BOTTOM = Inches(0.6)
    
    # 内容区
    CONTENT_WIDTH = WIDTH - MARGIN_LEFT - MARGIN_RIGHT
    CONTENT_HEIGHT = HEIGHT - MARGIN_TOP - MARGIN_BOTTOM


# ══════════════════════════════════════════════════════════════
# 主题配色
# ══════════════════════════════════════════════════════════════

THEMES = {
    "executive": {
        "bg": RGBColor(13, 27, 42),        # 深蓝背景
        "card_bg": RGBColor(27, 40, 56),    # 卡片背景
        "text": RGBColor(255, 255, 255),    # 白色文字
        "muted": RGBColor(136, 153, 170),   # 灰色副文字
        "accent": RGBColor(255, 183, 0),    # 金色强调
        "light_accent": RGBColor(255, 220, 100),  # 浅金色
        "success": RGBColor(0, 230, 118),   # 绿色（成功）
        "warning": RGBColor(255, 152, 0),   # 橙色（警告）
        "danger": RGBColor(244, 67, 54),    # 红色（危险）
        "border": RGBColor(45, 63, 82),     # 边框色
    },
    "spark": {
        "bg": RGBColor(15, 12, 41),
        "card_bg": RGBColor(30, 25, 60),
        "text": RGBColor(240, 240, 240),
        "muted": RGBColor(160, 176, 192),
        "accent": RGBColor(36, 198, 220),
        "light_accent": RGBColor(100, 220, 240),
        "border": RGBColor(60, 50, 100),
    },
    "dark": {
        "bg": RGBColor(26, 26, 26),
        "card_bg": RGBColor(40, 40, 40),
        "text": RGBColor(255, 255, 255),
        "muted": RGBColor(150, 150, 150),
        "accent": RGBColor(0, 200, 100),
        "light_accent": RGBColor(100, 255, 150),
        "border": RGBColor(60, 60, 60),
    },
}

DEFAULT_THEME = "executive"


class PPTGeneratorV3:
    """V3.0 PPT生成器 - 美学升级版"""
    
    def __init__(self, theme=DEFAULT_THEME):
        self.prs = Presentation()
        self.prs.slide_width = Layout.WIDTH
        self.prs.slide_height = Layout.HEIGHT
        self.t = THEMES.get(theme, THEMES[DEFAULT_THEME])
    
    # ══════════════════════════════════════════════════════════
    # 核心方法：建立视觉层次
    # ══════════════════════════════════════════════════════════
    
    def _add_bg(self, slide):
        """添加背景"""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Layout.WIDTH, Layout.HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.t['bg']
        bg.line.fill.background()
        # 移到最底层
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)
    
    def _title_hero(self, slide, text, size=72, color=None):
        """超大字标题 - 核心视觉元素"""
        box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, 
            Inches(2.0), 
            Layout.CONTENT_WIDTH, 
            Inches(3)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = color or self.t['text']
        p.alignment = PP_ALIGN.LEFT
        return box
    
    def _subtitle_small(self, slide, text, top=Inches(4.5)):
        """小字副标题"""
        box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT,
            top,
            Layout.CONTENT_WIDTH,
            Inches(1)
        )
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.color.rgb = self.t['muted']
        p.alignment = PP_ALIGN.LEFT
        return box
    
    def _accent_line(self, slide, top=Inches(1.8)):
        """强调线 - 视觉引导"""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Layout.MARGIN_LEFT, top,
            Inches(1.5), Inches(0.08)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.t['accent']
        line.line.fill.background()
        return line
    
    def _card(self, slide, left, top, width, height):
        """卡片容器"""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.t['card_bg']
        card.line.fill.background()
        return card
    
    # ══════════════════════════════════════════════════════════
    # 页面类型
    # ══════════════════════════════════════════════════════════
    
    def page_title(self, main, sub=None, tag=None):
        """🎯 标题页 - 大留白，大标题"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 左侧强调线
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Layout.MARGIN_LEFT, Inches(2.5),
            Inches(0.12), Inches(2.5)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = t['accent']
        accent.line.fill.background()
        
        # 主标题 - 超大！
        title_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(2.3), 
            Layout.CONTENT_WIDTH - Inches(0.4), Inches(2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = main
        p.font.size = Pt(66)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 副标题
        if sub:
            sub_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(4.5),
                Layout.CONTENT_WIDTH, Inches(1)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = sub
            p.font.size = Pt(28)
            p.font.color.rgb = t['muted']
        
        # 标签（右下角）
        if tag:
            tag_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(6.2),
                Layout.CONTENT_WIDTH, Inches(0.5)
            )
            tf = tag_box.text_frame
            p = tf.paragraphs[0]
            p.text = tag
            p.font.size = Pt(16)
            p.font.color.rgb = t['accent']
        
        return slide
    
    def page_big_number(self, number, label, context=None):
        """📊 大数据页 - 数字震撼"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 数字 - 超级大！
        num_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Inches(1.5),
            Layout.CONTENT_WIDTH, Inches(3.5)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(160)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        
        # 标签 - 中等
        label_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Inches(5.2),
            Layout.CONTENT_WIDTH, Inches(0.8)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(36)
        p.font.color.rgb = t['text']
        
        # 上下文 - 小
        if context:
            ctx_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(6.2),
                Layout.CONTENT_WIDTH, Inches(0.5)
            )
            tf = ctx_box.text_frame
            p = tf.paragraphs[0]
            p.text = context
            p.font.size = Pt(20)
            p.font.color.rgb = t['muted']
        
        return slide
    
    def page_two_columns(self, left_title, left_items, right_title, right_items, page_title=None):
        """📐 对比页 - 左右分栏"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 页面标题
        if page_title:
            title_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
                Layout.CONTENT_WIDTH, Inches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = page_title
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = t['text']
        
        content_top = Inches(1.8) if page_title else Inches(1.2)
        
        # 左栏
        left_width = (Layout.CONTENT_WIDTH - Inches(0.4)) / 2
        
        # 左栏标题
        lt_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, content_top,
            left_width, Inches(0.6)
        )
        tf = lt_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        
        # 左栏内容 - 大留白
        lc_top = content_top + Inches(0.8)
        lc_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, lc_top,
            left_width, Inches(4.5)
        )
        tf = lc_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = t['text']
            p.space_before = Pt(18)
        
        # 中间分隔线
        mid_x = Layout.MARGIN_LEFT + left_width + Inches(0.2)
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, mid_x, content_top, Inches(0.02), Inches(5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = t['border']
        divider.line.fill.background()
        
        # 右栏
        right_left = mid_x + Inches(0.2)
        
        rt_box = slide.shapes.add_textbox(
            right_left, content_top,
            left_width, Inches(0.6)
        )
        tf = rt_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        
        rc_box = slide.shapes.add_textbox(
            right_left, lc_top,
            left_width, Inches(4.5)
        )
        tf = rc_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = t['text']
            p.space_before = Pt(18)
        
        return slide
    
    def page_chart_visual(self, data, chart_type='radar'):
        """📈 图表页 - 图表+文字配合，左右分栏"""
        from chart_generator import (
            create_radar_chart, create_funnel_chart,
            create_matrix_chart, create_pyramid_chart,
            create_gantt_chart, create_heatmap_chart,
            create_gauge_chart, create_wordcloud_chart,
            create_combo_chart, create_stacked_bar_chart
        )
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        title = data.get('title', '')
        subtitle = data.get('subtitle', '')
        description = data.get('description', '')
        
        # ===== 左侧：图表区域（约60%）=====
        chart_left = Layout.MARGIN_LEFT
        chart_top = Inches(1.4)
        chart_width = Inches(7.5)
        chart_height = Inches(5.8)
        
        # 生成图表
        try:
            chart_img_path = f'/tmp/v3_{chart_type}.png'
            
            if chart_type == 'radar':
                items = data.get('items', [])
                values = [float(item.get('value', 0)) for item in items]
                labels = [f"D{i+1}" for i in range(len(items))]
                img = create_radar_chart(values, labels, "")
                img.save(chart_img_path)
                
            elif chart_type == 'funnel':
                levels = data.get('levels', [])
                values = [float(l.get('value', 0)) for l in levels]
                labels = [l.get('title', '') for l in levels]
                img = create_funnel_chart(values, labels, "")
                img.save(chart_img_path)
                
            elif chart_type == 'matrix':
                quads = data.get('quadrants', [])
                img = create_matrix_chart(quads, "", "", "")
                img.save(chart_img_path)
                
            elif chart_type == 'pyramid':
                levels = data.get('levels', [])
                values = [int(l.get('title', '30').split('%')[0].strip()) if '%' in l.get('title', '') else 30 for l in levels]
                labels = [l.get('title', '') for l in levels]
                img = create_pyramid_chart(values, labels, "")
                img.save(chart_img_path)
                
            elif chart_type == 'gantt':
                tasks = data.get('tasks', [])
                img = create_gantt_chart(tasks, "")
                img.save(chart_img_path)
                
            elif chart_type == 'heatmap':
                matrix = data.get('matrix', [])
                xlabels = data.get('xlabels', [])
                ylabels = data.get('ylabels', [])
                if matrix and len(matrix) > 0:
                    import numpy as np
                    data_arr = np.array(matrix)
                    img = create_heatmap_chart(data_arr, xlabels, ylabels, "")
                    img.save(chart_img_path)
                
            elif chart_type == 'gauge':
                value = float(data.get('value', 0))
                label = data.get('label', '')
                min_val = float(data.get('min', 0))
                max_val = float(data.get('max', 100))
                img = create_gauge_chart(value, label, min_val, max_val, "")
                img.save(chart_img_path)
                
            elif chart_type == 'wordcloud':
                words = data.get('words', {})
                img = create_wordcloud_chart(words, "")
                img.save(chart_img_path)
                
            elif chart_type == 'combo':
                categories = data.get('categories', [])
                bar_values = data.get('bar_values', [])
                line_values = data.get('line_values', [])
                img = create_combo_chart(categories, bar_values, line_values, "")
                img.save(chart_img_path)
                
            elif chart_type == 'stacked':
                categories = data.get('categories', [])
                segments = data.get('segments', {})
                img = create_stacked_bar_chart(categories, segments, "")
                img.save(chart_img_path)
            
            # 插入图表
            slide.shapes.add_picture(chart_img_path, chart_left, chart_top, width=chart_width)
            
        except Exception as e:
            print(f"图表生成失败: {e}")
            import traceback
            traceback.print_exc()
        
        # ===== 右侧：文字说明区域（约40%）=====
        right_left = Inches(8.5)
        right_top = Inches(1.4)
        right_width = Inches(4.5)
        
        # 页面大标题
        if title:
            title_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
                Layout.CONTENT_WIDTH, Inches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = t['text']
        
        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(1.0),
                Layout.CONTENT_WIDTH, Inches(0.4)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = t['muted']
        
        # 右侧标题
        right_title_box = slide.shapes.add_textbox(
            right_left, right_top,
            right_width, Inches(0.5)
        )
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "📊 核心洞察"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        
        # 右侧分隔线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            right_left, right_top + Inches(0.55),
            right_width, Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = t['border']
        line.line.fill.background()
        
        # 右侧内容 - 根据图表类型生成不同说明
        content_y = right_top + Inches(0.7)
        
        if chart_type == 'radar':
            self._render_radar_insights(slide, data, right_left, content_y, right_width)
        elif chart_type == 'funnel':
            self._render_funnel_insights(slide, data, right_left, content_y, right_width)
        elif chart_type == 'gantt':
            self._render_gantt_insights(slide, data, right_left, content_y, right_width)
        elif chart_type == 'heatmap':
            self._render_heatmap_insights(slide, data, right_left, content_y, right_width)
        elif chart_type == 'matrix':
            self._render_matrix_insights(slide, data, right_left, content_y, right_width)
        elif chart_type == 'pyramid':
            self._render_pyramid_insights(slide, data, right_left, content_y, right_width)
        elif chart_type == 'combo':
            self._render_combo_insights(slide, data, right_left, content_y, right_width)
        elif chart_type == 'stacked':
            self._render_stacked_insights(slide, data, right_left, content_y, right_width)
        elif chart_type == 'gauge':
            self._render_gauge_insights(slide, data, right_left, content_y, right_width)
        else:
            # 默认渲染items
            self._render_items_legend_v2(slide, data, right_left, content_y, right_width)
        
        return slide
    
    def _render_radar_insights(self, slide, data, left, top, width):
        """雷达图 - 详细洞察"""
        t = self.t
        items = data.get('items', [])
        
        # 标题
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "能力评估分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 计算洞察
        values = [float(item.get('value', 0)) for item in items]
        avg = sum(values) / len(values) if values else 0
        max_val = max(values) if values else 0
        min_val = min(values) if values else 0
        max_item = max(items, key=lambda x: float(x.get('value', 0))) if items else {}
        min_item = min(items, key=lambda x: float(x.get('value', 0))) if items else {}
        
        y_offset = top + Inches(0.5)
        
        # 平均值
        avg_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
        tf = avg_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• 综合得分：{avg:.0f}分"
        p.font.size = Pt(14)
        p.font.color.rgb = t['text']
        y_offset += Inches(0.4)
        
        # 最高项
        if max_item:
            max_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = max_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• 优势能力：{max_item.get('title', '')}"
            p.font.size = Pt(14)
            p.font.color.rgb = t['success']
            p2 = tf.add_paragraph()
            p2.text = f"  {max_item.get('value', 0)}分，领先平均水平{max_val - avg:.0f}分"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
        
        # 最低项
        if min_item:
            min_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = min_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• 待提升：{min_item.get('title', '')}"
            p.font.size = Pt(14)
            p.font.color.rgb = t['warning']
            p2 = tf.add_paragraph()
            p2.text = f"  {min_item.get('value', 0)}分，需提升{min_val + 10 - float(min_item.get('value', 0)):.0f}分达平均"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
        
        # 建议
        suggest_box = slide.shapes.add_textbox(left, y_offset + Inches(0.3), width, Inches(0.6))
        tf = suggest_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "💡 建议"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = t['accent']
        p2 = tf.add_paragraph()
        p2.text = f"聚焦提升{min_item.get('title', '')}能力，预计可提升综合得分{avg + 5:.0f}分"
        p2.font.size = Pt(12)
        p2.font.color.rgb = t['muted']
    
    def _render_funnel_insights(self, slide, data, left, top, width):
        """漏斗图 - 详细洞察"""
        t = self.t
        levels = data.get('levels', [])
        
        # 标题
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "转化漏斗分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        # 计算转化率
        if len(levels) >= 2:
            first_val = float(levels[0].get('value', 1))
            last_val = float(levels[-1].get('value', 0))
            total_conv = last_val / first_val * 100 if first_val > 0 else 0
            
            conv_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = conv_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• 总体转化率：{total_conv:.2f}%"
            p.font.size = Pt(14)
            p.font.color.rgb = t['text']
            y_offset += Inches(0.4)
            
            # 各阶段转化
            for i in range(len(levels) - 1):
                curr = float(levels[i].get('value', 0))
                next_val = float(levels[i+1].get('value', 0))
                conv = next_val / curr * 100 if curr > 0 else 0
                
                stage_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.35))
                tf = stage_box.text_frame
                p = tf.paragraphs[0]
                curr_name = levels[i].get('title', '').split()[0]
                next_name = levels[i+1].get('title', '').split()[0]
                p.text = f"• {curr_name}→{next_name}"
                p.font.size = Pt(13)
                p.font.color.rgb = t['muted']
                p2 = tf.add_paragraph()
                p2.text = f"  转化率 {conv:.1f}%，流失 {curr - next_val:,.0f}人"
                p2.font.size = Pt(12)
                p2.font.color.rgb = t['muted']
                y_offset += Inches(0.5)
        
        # 关键流失点
        if len(levels) >= 2:
            losses = []
            for i in range(len(levels) - 1):
                curr = float(levels[i].get('value', 0))
                next_val = float(levels[i+1].get('value', 0))
                loss = curr - next_val
                losses.append((levels[i].get('title', ''), loss))
            
            max_loss = max(losses, key=lambda x: x[1]) if losses else ('', 0)
            
            loss_box = slide.shapes.add_textbox(left, y_offset + Inches(0.2), width, Inches(0.5))
            tf = loss_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "⚠️ 最大流失点"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = t['danger']
            p2 = tf.add_paragraph()
            p2.text = f"{max_loss[0]}环节流失最多，建议重点优化"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
    
    def _render_gantt_insights(self, slide, data, left, top, width):
        """甘特图 - 详细洞察"""
        t = self.t
        tasks = data.get('tasks', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "项目进度分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        # 统计
        total = len(tasks)
        complete = len([t for t in tasks if t.get('status') == 'complete'])
        in_progress = len([t for t in tasks if t.get('status') == 'in_progress'])
        pending = len([t for t in tasks if t.get('status') == 'pending'])
        
        stats = f"• 总任务：{total}个"
        stats_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.35))
        tf = stats_box.text_frame
        p = tf.paragraphs[0]
        p.text = stats
        p.font.size = Pt(13)
        p.font.color.rgb = t['text']
        y_offset += Inches(0.35)
        
        status_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.9))
        tf = status_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"✅ 已完成：{complete}个"
        p.font.size = Pt(13)
        p.font.color.rgb = t['success']
        p2 = tf.add_paragraph()
        p2.text = f"🔄 进行中：{in_progress}个"
        p2.font.size = Pt(13)
        p2.font.color.rgb = t['accent']
        p3 = tf.add_paragraph()
        p3.text = f"⏳ 待开始：{pending}个"
        p3.font.size = Pt(13)
        p3.font.color.rgb = t['muted']
        y_offset += Inches(0.9)
        
        # 关键里程碑
        in_progress_tasks = [t for t in tasks if t.get('status') == 'in_progress']
        if in_progress_tasks:
            milestone_box = slide.shapes.add_textbox(left, y_offset + Inches(0.2), width, Inches(0.5))
            tf = milestone_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "🎯 当前重点"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = t['accent']
            for task in in_progress_tasks[:2]:
                p2 = tf.add_paragraph()
                p2.text = f"• {task.get('name', '')}（{task.get('progress', 0)}%）"
                p2.font.size = Pt(12)
                p2.font.color.rgb = t['muted']
    
    def _render_heatmap_insights(self, slide, data, left, top, width):
        """热力图 - 详细洞察"""
        t = self.t
        matrix = data.get('matrix', [])
        xlabels = data.get('xlabels', [])
        ylabels = data.get('ylabels', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "区域表现分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        if matrix and len(matrix) > 0 and len(matrix[0]) > 0:
            import numpy as np
            data_arr = np.array(matrix)
            
            # 最高值
            max_idx = np.unravel_index(data_arr.argmax(), data_arr.shape)
            max_val = data_arr.max()
            max_row = ylabels[max_idx[0]] if max_idx[0] < len(ylabels) else ''
            max_col = xlabels[max_idx[1]] if max_idx[1] < len(xlabels) else ''
            
            max_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = max_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"🏆 最佳表现"
            p.font.size = Pt(14)
            p.font.color.rgb = t['success']
            p2 = tf.add_paragraph()
            p2.text = f"{max_row} × {max_col} = {max_val:.0f}"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
            
            # 最低值
            min_idx = np.unravel_index(data_arr.argmin(), data_arr.shape)
            min_val = data_arr.min()
            min_row = ylabels[min_idx[0]] if min_idx[0] < len(ylabels) else ''
            min_col = xlabels[min_idx[1]] if min_idx[1] < len(xlabels) else ''
            
            min_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = min_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"📍 待提升区域"
            p.font.size = Pt(14)
            p.font.color.rgb = t['warning']
            p2 = tf.add_paragraph()
            p2.text = f"{min_row} × {min_col} = {min_val:.0f}，需重点关注"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
    
    def _render_matrix_insights(self, slide, data, left, top, width):
        """矩阵图 - 详细洞察"""
        t = self.t
        quads = data.get('quadrants', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "产品组合分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        # 各象限分析
        quadrant_names = ["明星产品", "问题产品", "瘦狗产品", "金牛产品"]
        colors = [t['success'], t['warning'], t['danger'], t['accent']]
        
        for i, (quadrant, name, color) in enumerate(zip(quads, quadrant_names, colors)):
            items = quadrant if isinstance(quadrant, list) else (quadrant.get('items', []) if isinstance(quadrant, dict) else [])
            if not items:
                continue
                
            q_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.35))
            tf = q_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {name}：{len(items)}个"
            p.font.size = Pt(13)
            p.font.color.rgb = color
            y_offset += Inches(0.35)
            
            if i == 0:  # 明星产品
                p2 = tf.add_paragraph()
                p2.text = "  高增长高份额，重点投入资源"
                p2.font.size = Pt(11)
                p2.font.color.rgb = t['muted']
            elif i == 3:  # 金牛产品
                p2 = tf.add_paragraph()
                p2.text = "  稳定盈利，贡献现金流"
                p2.font.size = Pt(11)
                p2.font.color.rgb = t['muted']
            y_offset += Inches(0.15)
    
    def _render_pyramid_insights(self, slide, data, left, top, width):
        """金字塔 - 详细洞察"""
        t = self.t
        levels = data.get('levels', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "品类结构分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        # 解读
        for level in levels[:3]:
            level_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.6))
            tf = level_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            title = level.get('title', '')
            desc = level.get('desc', '')
            p.text = f"• {title}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = t['accent']
            if desc:
                p2 = tf.add_paragraph()
                p2.text = f"  {desc}"
                p2.font.size = Pt(11)
                p2.font.color.rgb = t['muted']
            y_offset += Inches(0.6)
    
    def _render_combo_insights(self, slide, data, left, top, width):
        """组合图 - 详细洞察"""
        t = self.t
        bar_values = data.get('bar_values', [])
        line_values = data.get('line_values', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "趋势分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        if bar_values and line_values:
            import numpy as np
            bar_arr = np.array(bar_values)
            line_arr = np.array(line_values)
            
            # 趋势
            trend = "上升" if line_arr[-1] > line_arr[0] else "下降"
            change = line_arr[-1] - line_arr[0]
            
            trend_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = trend_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• 增长率趋势：{trend}"
            p.font.size = Pt(13)
            p.font.color.rgb = t['text']
            p2 = tf.add_paragraph()
            p2.text = f"  {change:+.1f}个百分点"
            p2.font.size = Pt(12)
            p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
            
            # 峰值
            max_bar_idx = np.argmax(bar_arr)
            max_bar = bar_arr.max()
            
            peak_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
            tf = peak_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• 峰值：{max_bar:.0f}"
            p.font.size = Pt(13)
            p.font.color.rgb = t['success']
            p2 = tf.add_paragraph()
            categories = data.get('categories', [])
            if max_bar_idx < len(categories):
                p2.text = f"  出现在{categories[max_bar_idx]}"
                p2.font.size = Pt(12)
                p2.font.color.rgb = t['muted']
    
    def _render_stacked_insights(self, slide, data, left, top, width):
        """堆叠图 - 详细洞察"""
        t = self.t
        segments = data.get('segments', {})
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "构成变化分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        seg_names = list(segments.keys())
        if len(seg_names) >= 2:
            first_vals = [segments[seg][0] if len(segments[seg]) > 0 else 0 for seg in seg_names]
            last_vals = [segments[seg][-1] if len(segments[seg]) > 0 else 0 for seg in seg_names]
            
            changes = [last - first for first, last in zip(first_vals, last_vals)]
            
            # 增长最快的
            max_change = max(changes)
            max_idx = changes.index(max_change)
            
            change_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.5))
            tf = change_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"📈 增长最快：{seg_names[max_idx]}"
            p.font.size = Pt(13)
            p.font.color.rgb = t['success']
            p2 = tf.add_paragraph()
            p2.text = f"  占比从{first_vals[max_idx]}%增至{last_vals[max_idx]}%（+{max_change}%）"
            p2.font.size = Pt(11)
            p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
            
            # 下降的
            min_change = min(changes)
            if min_change < 0:
                min_idx = changes.index(min_change)
                decline_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.5))
                tf = decline_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"📉 占比下降：{seg_names[min_idx]}"
                p.font.size = Pt(13)
                p.font.color.rgb = t['danger']
                p2 = tf.add_paragraph()
                p2.text = f"  占比从{first_vals[min_idx]}%降至{last_vals[min_idx]}%（{min_change}%）"
                p2.font.size = Pt(11)
                p2.font.color.rgb = t['muted']
    
    def _render_gauge_insights(self, slide, data, left, top, width):
        """仪表盘 - 详细洞察"""
        t = self.t
        value = float(data.get('value', 0))
        label = data.get('label', '')
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "KPI达成分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        # 状态评估
        if value >= 80:
            status = "优秀"
            color = t['success']
        elif value >= 60:
            status = "良好"
            color = t['accent']
        elif value >= 40:
            status = "一般"
            color = t['warning']
        else:
            status = "需改进"
            color = t['danger']
        
        status_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.4))
        tf = status_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"• 状态：{status}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        y_offset += Inches(0.4)
        
        # 解读
        gap = 100 - value
        gap_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.5))
        tf = gap_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"• 距目标差距"
        p.font.size = Pt(13)
        p.font.color.rgb = t['text']
        p2 = tf.add_paragraph()
        p2.text = f"还差{gap:.0f}个百分点达到100%目标"
        p2.font.size = Pt(12)
        p2.font.color.rgb = t['muted']
    
    def _render_items_legend_v2(self, slide, data, left, top, width):
        """通用items渲染"""
        t = self.t
        items = data.get('items', [])
        
        title_box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "数据要点"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        y_offset = top + Inches(0.5)
        
        for item in items[:5]:
            item_box = slide.shapes.add_textbox(left, y_offset, width, Inches(0.5))
            tf = item_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"• {item.get('title', '')}"
            p.font.size = Pt(13)
            p.font.color.rgb = t['text']
            if item.get('value'):
                p2 = tf.add_paragraph()
                p2.text = f"  {item['value']}"
                p2.font.size = Pt(12)
                p2.font.color.rgb = t['muted']
            y_offset += Inches(0.5)
    
    def _render_items_legend(self, slide, items, max_val=100):
        """渲染指标图例"""
        t = self.t
        legend_x = Inches(8.8)
        legend_y = Inches(2.0)
        
        for i, item in enumerate(items):
            # 标签
            label_box = slide.shapes.add_textbox(
                legend_x, legend_y + Inches(i * 0.8),
                Inches(4), Inches(0.7)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {item.get('title', '')}"
            p.font.size = Pt(18)
            p.font.color.rgb = t['text']
            
            # 值
            val_box = slide.shapes.add_textbox(
                legend_x + Inches(0.3), legend_y + Inches(i * 0.8 + 0.35),
                Inches(3.5), Inches(0.4)
            )
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(item.get('value', ''))
            p.font.size = Pt(14)
            p.font.color.rgb = t['muted']
            
            # 进度条
            val = float(item.get('value', 0)) / max_val
            bar_width = val * 3.5
            
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                legend_x, legend_y + Inches(i * 0.8 + 0.55),
                Inches(bar_width), Inches(0.08)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = t['accent']
            bar.line.fill.background()
    
    def _render_levels_legend(self, slide, levels):
        """渲染层级图例"""
        t = self.t
        legend_x = Inches(9)
        legend_y = Inches(2.0)
        
        for i, level in enumerate(levels):
            label_box = slide.shapes.add_textbox(
                legend_x, legend_y + Inches(i * 1.0),
                Inches(4), Inches(0.9)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = level.get('title', '')
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = t['accent']
            
            if level.get('desc'):
                p2 = tf.add_paragraph()
                p2.text = level['desc']
                p2.font.size = Pt(14)
                p2.font.color.rgb = t['muted']
    
    def page_content_focus(self, title, items, style='cards'):
        """📝 内容页 - 聚焦要点"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 页面标题
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
            Layout.CONTENT_WIDTH, Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        if style == 'cards':
            # 卡片式布局
            card_width = (Layout.CONTENT_WIDTH - Inches(0.6)) / 2
            card_height = Inches(1.8)
            
            for i, item in enumerate(items):
                row = i // 2
                col = i % 2
                
                left = Layout.MARGIN_LEFT + col * (card_width + Inches(0.2))
                top = Inches(2.0) + row * (card_height + Inches(0.3))
                
                # 卡片背景
                card = self._card(slide, left, top, card_width, card_height)
                
                # 标题
                item_title = item.get('title', '') if isinstance(item, dict) else item
                title_tb = slide.shapes.add_textbox(
                    left + Inches(0.3), top + Inches(0.3),
                    card_width - Inches(0.6), Inches(0.6)
                )
                tf = title_tb.text_frame
                p = tf.paragraphs[0]
                p.text = item_title
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = t['accent']
                
                # 值
                if isinstance(item, dict) and item.get('value'):
                    val_tb = slide.shapes.add_textbox(
                        left + Inches(0.3), top + Inches(0.9),
                        card_width - Inches(0.6), Inches(0.5)
                    )
                    tf = val_tb.text_frame
                    p = tf.paragraphs[0]
                    p.text = item['value']
                    p.font.size = Pt(18)
                    p.font.color.rgb = t['text']
        
        return slide
    
    def page_quote(self, quote, attribution=None):
        """💬 引用页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 大引号
        quote_mark = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Inches(1.5),
            Inches(2), Inches(2)
        )
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = '"'
        p.font.size = Pt(200)
        p.font.color.rgb = t['accent']
        p.font.name = 'Georgia'
        
        # 引用内容
        quote_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.8),
            Layout.CONTENT_WIDTH - Inches(1), Inches(3)
        )
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = quote
        p.font.size = Pt(32)
        p.font.italic = True
        p.font.color.rgb = t['text']
        
        # 作者
        if attribution:
            attr_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(6),
                Layout.CONTENT_WIDTH, Inches(0.6)
            )
            tf = attr_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {attribution}"
            p.font.size = Pt(20)
            p.font.color.rgb = t['muted']
            p.alignment = PP_ALIGN.RIGHT
        
        return slide
    
    def page_timeline(self, title, steps):
        """📅 时间线页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Layout.MARGIN_TOP,
            Layout.CONTENT_WIDTH, Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = t['text']
        
        # 时间线
        num = len(steps)
        start_x = Layout.MARGIN_LEFT + Inches(1)
        end_x = Layout.WIDTH - Layout.MARGIN_RIGHT - Inches(1)
        line_y = Inches(4.2)
        
        # 主线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            start_x, line_y, end_x - start_x, Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = t['border']
        line.line.fill.background()
        
        # 节点和内容
        for i, step in enumerate(steps):
            x = start_x + (end_x - start_x) * i / max(num - 1, 1)
            
            # 节点圆
            node = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x - Inches(0.2), line_y - Inches(0.2),
                Inches(0.4), Inches(0.4)
            )
            node.fill.solid()
            node.fill.fore_color.rgb = t['accent']
            node.line.fill.background()
            
            # 步骤标题
            step_title = step.get('title', f'Step {i+1}')
            st_box = slide.shapes.add_textbox(
                x - Inches(1), line_y + Inches(0.5),
                Inches(2), Inches(0.6)
            )
            tf = st_box.text_frame
            p = tf.paragraphs[0]
            p.text = step_title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = t['text']
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            if step.get('desc'):
                desc_box = slide.shapes.add_textbox(
                    x - Inches(1.2), line_y + Inches(1.1),
                    Inches(2.4), Inches(1)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = step['desc']
                p.font.size = Pt(14)
                p.font.color.rgb = t['muted']
                p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def page_closing(self, title, cta=None, contact=None):
        """🎬 结尾页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        t = self.t
        
        # 主标题 - 居中大标题
        title_box = slide.shapes.add_textbox(
            Layout.MARGIN_LEFT, Inches(2.5),
            Layout.CONTENT_WIDTH, Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = t['text']
        p.alignment = PP_ALIGN.CENTER
        
        # CTA
        if cta:
            cta_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(4.5),
                Layout.CONTENT_WIDTH, Inches(0.8)
            )
            tf = cta_box.text_frame
            p = tf.paragraphs[0]
            p.text = cta
            p.font.size = Pt(24)
            p.font.color.rgb = t['accent']
            p.alignment = PP_ALIGN.CENTER
        
        # 联系
        if contact:
            contact_box = slide.shapes.add_textbox(
                Layout.MARGIN_LEFT, Inches(6),
                Layout.CONTENT_WIDTH, Inches(0.5)
            )
            tf = contact_box.text_frame
            p = tf.paragraphs[0]
            p.text = contact
            p.font.size = Pt(16)
            p.font.color.rgb = t['muted']
            p.alignment = PP_ALIGN.CENTER
        
        return slide


def create_ppt_v3(slides, output_path, theme='executive'):
    """V3.0创建PPT"""
    gen = PPTGeneratorV3(theme)
    
    for slide in slides:
        stype = slide.get('type', '')
        
        if stype == 'title':
            gen.page_title(
                slide.get('main', ''),
                slide.get('sub'),
                slide.get('tag')
            )
        elif stype == 'big_number':
            gen.page_big_number(
                slide.get('number', ''),
                slide.get('label', ''),
                slide.get('context')
            )
        elif stype == 'two_columns':
            gen.page_two_columns(
                slide.get('left_title', ''),
                slide.get('left_items', []),
                slide.get('right_title', ''),
                slide.get('right_items', []),
                slide.get('page_title')
            )
        elif stype == 'radar':
            gen.page_chart_visual(slide, 'radar')
        elif stype == 'funnel':
            gen.page_chart_visual(slide, 'funnel')
        elif stype == 'matrix':
            gen.page_chart_visual(slide, 'matrix')
        elif stype == 'pyramid':
            gen.page_chart_visual(slide, 'pyramid')
        elif stype == 'gantt':
            gen.page_chart_visual(slide, 'gantt')
        elif stype == 'heatmap':
            gen.page_chart_visual(slide, 'heatmap')
        elif stype == 'gauge':
            gen.page_chart_visual(slide, 'gauge')
        elif stype == 'wordcloud':
            gen.page_chart_visual(slide, 'wordcloud')
        elif stype == 'combo':
            gen.page_chart_visual(slide, 'combo')
        elif stype == 'stacked':
            gen.page_chart_visual(slide, 'stacked')
        elif stype == 'content_cards':
            gen.page_content_focus(
                slide.get('title', ''),
                slide.get('items', []),
                'cards'
            )
        elif stype == 'quote':
            gen.page_quote(
                slide.get('quote', ''),
                slide.get('attribution')
            )
        elif stype == 'timeline':
            gen.page_timeline(
                slide.get('title', ''),
                slide.get('steps', [])
            )
        elif stype == 'closing':
            gen.page_closing(
                slide.get('title', ''),
                slide.get('cta'),
                slide.get('contact')
            )
    
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    gen.prs.save(output_path)
    return output_path
